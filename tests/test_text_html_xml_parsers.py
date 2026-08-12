# -*- coding: utf-8 -*-
"""TXT/HTML/XML 解析器测试（Phase 2 Task 5 剩余）。"""
from __future__ import annotations

from pathlib import Path

import pytest

from src.data_prep.artifact_store import ArtifactStore
from src.data_prep.models import RawArtifact
from src.data_prep.document_models import DocumentElement
from src.parsers.text_html_xml import HtmlParser, TextParser, XmlParser
from src.parsers.markdown import MarkdownParser
from src.parsers.presentation import PresentationParser


def _make_artifact(store: ArtifactStore, task_id: str, data: bytes, ext: str, media_type: str) -> RawArtifact:
    return store.write_raw(
        task_id=task_id, source_id="file-1", data=data,
        uri=f"upload://data.{ext}", media_type=media_type, ext=ext,
    )


def test_txt_parser_parses_lines(tmp_path: Path):
    store = ArtifactStore(root=str(tmp_path))
    art = _make_artifact(store, "task-1", b"line one\nline two\nline three\n", "txt", "text/plain")

    parser = TextParser()
    records, rejects = parser.parse(art, b"line one\nline two\nline three\n")

    assert len(records) == 3
    assert records[0].data["text"] == "line one"
    assert records[0].meta["position"]["line"] == 1
    assert records[0].meta["parser"] == "txt"
    DocumentElement.model_validate(records[0].data["elements"][0])
    assert rejects == []


def test_txt_parser_handles_gbk(tmp_path: Path):
    store = ArtifactStore(root=str(tmp_path))
    data = "第一行\n第二行\n".encode("gbk")
    art = _make_artifact(store, "task-1", data, "txt", "text/plain")

    parser = TextParser()
    records, _ = parser.parse(art, data)

    assert len(records) == 2
    assert records[0].data["text"] == "第一行"


def test_html_parser_extracts_table_rows(tmp_path: Path):
    store = ArtifactStore(root=str(tmp_path))
    data = b"<table><tr><th>id</th><th>name</th></tr><tr><td>1</td><td>Alice</td></tr><tr><td>2</td><td>Bob</td></tr></table>"
    art = _make_artifact(store, "task-1", data, "html", "text/html")

    parser = HtmlParser()
    records, rejects = parser.parse(art, data)

    assert len(records) == 2
    assert records[0].data["id"] == "1"
    assert records[0].data["name"] == "Alice"
    assert records[0].meta["parser"] == "html_table"
    assert rejects == []


def test_html_parser_extracts_paragraphs_when_no_table(tmp_path: Path):
    store = ArtifactStore(root=str(tmp_path))
    data = b"<html><body><p>First paragraph</p><p>Second paragraph</p></body></html>"
    art = _make_artifact(store, "task-1", data, "html", "text/html")

    parser = HtmlParser()
    records, _ = parser.parse(art, data)

    assert len(records) == 2
    assert records[0].data["text"] == "First paragraph"
    assert records[0].meta["parser"] == "html_text"


def test_html_parser_keeps_text_and_table_in_dom_order(tmp_path: Path):
    store = ArtifactStore(root=str(tmp_path))
    data = (
        "<h1>付款条款</h1><p>验收后付款。</p>"
        "<table><tr><th>项目</th><th>比例</th></tr>"
        "<tr><td>首付款</td><td>60%</td></tr></table>"
        "<p>合同结束。</p>"
    ).encode("utf-8")
    artifact = _make_artifact(store, "task-html", data, "html", "text/html")

    records, rejects = HtmlParser().parse(artifact, data)
    elements = [
        DocumentElement.model_validate(record.data["elements"][0])
        for record in records
    ]

    assert rejects == []
    assert [item.text for item in elements] == [
        "付款条款",
        "验收后付款。",
        "项目：首付款；比例：60%",
        "合同结束。",
    ]
    assert [item.reading_order for item in elements] == [0, 1, 2, 3]


def test_xml_parser_extracts_child_nodes(tmp_path: Path):
    store = ArtifactStore(root=str(tmp_path))
    data = b'<root><item><id>1</id><name>Alice</name></item><item><id>2</id><name>Bob</name></item></root>'
    art = _make_artifact(store, "task-1", data, "xml", "application/xml")

    parser = XmlParser()
    records, rejects = parser.parse(art, data)

    assert len(records) == 2
    assert records[0].data["id"] == "1"
    assert records[0].data["name"] == "Alice"
    assert records[0].meta["parser"] == "xml"
    DocumentElement.model_validate(records[0].data["elements"][0])


def test_markdown_parser_preserves_heading_and_paragraph(tmp_path: Path):
    store = ArtifactStore(root=str(tmp_path))
    data = "# 付款条款\n\n验收后十五日内支付。\n".encode("utf-8")
    artifact = _make_artifact(
        store, "task-md", data, "md", "text/markdown"
    )

    records, rejects = MarkdownParser().parse(artifact, data)
    elements = [
        DocumentElement.model_validate(record.data["elements"][0])
        for record in records
    ]

    assert rejects == []
    assert [item.element_type.value for item in elements] == [
        "heading",
        "paragraph",
    ]
    assert [item.text for item in elements] == [
        "付款条款",
        "验收后十五日内支付。",
    ]


def test_markdown_parser_preserves_table_rows(tmp_path: Path):
    store = ArtifactStore(root=str(tmp_path))
    data = (
        "| 项目 | 比例 |\n| --- | --- |\n| 首付款 | 60% |\n"
    ).encode("utf-8")
    artifact = _make_artifact(
        store, "task-md-table", data, "md", "text/markdown"
    )

    records, rejects = MarkdownParser().parse(artifact, data)
    element = DocumentElement.model_validate(records[0].data["elements"][0])

    assert rejects == []
    assert element.element_type.value == "table"
    assert element.metadata["table_row"] == {"项目": "首付款", "比例": "60%"}


def test_pptx_parser_keeps_slide_shape_location(tmp_path: Path):
    import io
    from pptx import Presentation

    store = ArtifactStore(root=str(tmp_path))
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "付款条款"
    slide.placeholders[1].text = "验收后十五日内支付。"
    buffer = io.BytesIO()
    presentation.save(buffer)
    data = buffer.getvalue()
    artifact = _make_artifact(
        store,
        "task-pptx",
        data,
        "pptx",
        (
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ),
    )

    records, rejects = PresentationParser().parse(artifact, data)
    elements = [
        DocumentElement.model_validate(record.data["elements"][0])
        for record in records
    ]

    assert rejects == []
    assert [item.text for item in elements] == [
        "付款条款",
        "验收后十五日内支付。",
    ]
    assert all(item.page == 1 for item in elements)
    assert all(item.metadata["location"]["slide"] == 1 for item in elements)
    assert rejects == []
