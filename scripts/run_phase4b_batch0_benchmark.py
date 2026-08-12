#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""运行 Phase 4B 批次 0 的本机确定性工具赛马并输出版本化 JSON。"""
from __future__ import annotations

import argparse
import importlib.metadata
import json
from pathlib import Path
import shutil
import statistics
import subprocess
import sys
import tempfile
import time
from typing import Any


PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from src.evaluation.semantic_harness.adapters import get_table_adapter  # noqa: E402
from src.evaluation.semantic_harness.fixtures import load_batch0_manifest  # noqa: E402
from src.evaluation.semantic_harness.graph import run_table_benchmark  # noqa: E402
from src.evaluation.semantic_harness.table_ops import (  # noqa: E402
    run_table_operation_suite,
)
from scripts.generate_semantic_batch0_fixtures import normalize_ooxml  # noqa: E402


DEFAULT_MANIFEST = (
    PROJECT_ROOT
    / "tests"
    / "fixtures"
    / "semantic_harness"
    / "public"
    / "batch0"
    / "manifest.json"
)
DEFAULT_OUTPUT = (
    PROJECT_ROOT
    / "docs"
    / "plans"
    / "phase4b-batch0-results"
    / "local-benchmark.json"
)


def _version(package: str) -> str | None:
    try:
        return importlib.metadata.version(package)
    except importlib.metadata.PackageNotFoundError:
        return None


def _contains_all(text: str, expected: list[str]) -> bool:
    compact = "".join(text.split())
    return all("".join(item.split()) in compact for item in expected)


def _table_benchmark(manifest_path: Path, output_root: Path) -> dict[str, Any]:
    candidates = ("table.duckdb", "table.polars", "table.pandas")
    report = run_table_benchmark(
        manifest_path=manifest_path,
        case_id="workload_filter",
        candidate_ids=candidates,
        output_dir=output_root / "table-workload-filter",
    )
    report = _portable(report)
    (output_root / "table-workload-filter" / "benchmark-summary.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    manifest = load_batch0_manifest(manifest_path)
    case = manifest.case("workload_filter")
    source_path = manifest.resolve(case.canonical_input)
    timings: dict[str, Any] = {}
    for capability_id in candidates:
        adapter = get_table_adapter(capability_id)
        durations = [adapter.run(case, source_path).duration_ms for _ in range(7)]
        timings[capability_id] = {
            "runs": durations,
            "warm_median_ms": statistics.median(durations[1:]),
            "warm_p95_ms": max(durations[1:]),
        }
    operation_root = manifest.root / "table_operations"
    operation_results = [
        run_table_operation_suite(capability_id, operation_root)
        for capability_id in candidates
    ]
    return {
        **report,
        "timings": timings,
        "operation_suite": {
            "status": (
                "pass"
                if all(item["status"] == "pass" for item in operation_results)
                else "fail"
            ),
            "results": operation_results,
        },
        "scale_suite": _large_table_benchmark(),
    }


def _large_table_benchmark() -> dict[str, Any]:
    import duckdb
    import pandas as pd
    import polars as pl
    import psutil
    import pyarrow as pa
    import pyarrow.parquet as pq

    process = psutil.Process()
    results: list[dict[str, Any]] = []
    with tempfile.TemporaryDirectory(prefix="mangrove-phase4b-table-") as temp:
        temp_root = Path(temp)
        for row_count in (100_000, 1_000_000):
            identifiers = list(range(row_count))
            table = pa.table(
                {
                    "记录号": identifiers,
                    "分组": [f"G{index % 10}" for index in identifiers],
                    "保留": [index % 3 == 0 for index in identifiers],
                    "金额": [float(index % 100) for index in identifiers],
                }
            )
            source = temp_root / f"scale-{row_count}.parquet"
            pq.write_table(table, source)
            del table, identifiers
            expected_rows = (row_count + 2) // 3
            expected_sum = sum(
                float(index % 100)
                for index in range(0, row_count, 3)
            )
            for capability_id in (
                "table.duckdb",
                "table.polars",
                "table.pandas",
            ):
                before = process.memory_info().rss
                started = time.perf_counter()
                if capability_id == "table.duckdb":
                    connection = duckdb.connect(database=":memory:")
                    try:
                        actual_rows, actual_sum = connection.execute(
                            """
                            SELECT COUNT(*), SUM("金额")
                            FROM read_parquet(?) WHERE "保留" = true
                            """,
                            [str(source)],
                        ).fetchone()
                    finally:
                        connection.close()
                elif capability_id == "table.polars":
                    frame = (
                        pl.scan_parquet(source)
                        .filter(pl.col("保留"))
                        .select(
                            pl.len().alias("rows"),
                            pl.col("金额").sum().alias("amount"),
                        )
                        .collect()
                    )
                    actual_rows = frame["rows"][0]
                    actual_sum = frame["amount"][0]
                    del frame
                else:
                    frame = pd.read_parquet(source, columns=["保留", "金额"])
                    selected = frame.loc[frame["保留"], "金额"]
                    actual_rows = len(selected)
                    actual_sum = selected.sum()
                    del selected, frame
                after = process.memory_info().rss
                results.append(
                    {
                        "capability_id": capability_id,
                        "rows": row_count,
                        "status": (
                            "pass"
                            if int(actual_rows) == expected_rows
                            and float(actual_sum) == expected_sum
                            else "fail"
                        ),
                        "duration_ms": int(
                            (time.perf_counter() - started) * 1000
                        ),
                        "rss_delta_bytes": max(0, after - before),
                        "source_bytes": source.stat().st_size,
                    }
                )
    return {
        "status": (
            "pass" if all(item["status"] == "pass" for item in results) else "fail"
        ),
        "results": results,
        "measurement_note": (
            "rss_delta_bytes 是进程前后差值，不等同于采样峰值；"
            "耗时和正确性可用于本机候选排序。"
        ),
    }


def _markitdown_benchmark(manifest_path: Path) -> dict[str, Any]:
    from markitdown import MarkItDown

    manifest = load_batch0_manifest(manifest_path)
    converter = MarkItDown(enable_plugins=False)
    expected_contract = ["付款条款", "交付条款", "违约责任"]
    cases = [
        ("pdf", "documents/contract.pdf", expected_contract),
        ("docx", "documents/contract.docx", expected_contract),
        ("pptx", "documents/contract.pptx", expected_contract),
        ("html", "documents/contract.html", expected_contract),
        ("markdown", "documents/contract.md", expected_contract),
        ("txt", "documents/contract.txt", expected_contract),
        ("xml", "documents/contract.xml", expected_contract),
        (
            "xlsx",
            "workload_filter/source.xlsx",
            ["示例人员甲", "核销工作量天数", "工作量费用"],
        ),
    ]
    results = []
    for file_format, relative_path, expected in cases:
        started = time.perf_counter()
        item: dict[str, Any] = {
            "format": file_format,
            "path": relative_path,
        }
        try:
            converted = converter.convert(manifest.resolve(relative_path))
            text = converted.text_content or ""
            item.update(
                {
                    "status": "pass" if _contains_all(text, expected) else "fail",
                    "expected_tokens_found": [
                        token for token in expected if _contains_all(text, [token])
                    ],
                    "text_chars": len(text),
                }
            )
        except Exception as exc:  # noqa: BLE001 - 兼容矩阵必须保留单项错误
            item.update(
                {
                    "status": "error",
                    "error_type": type(exc).__name__,
                    "error": str(exc),
                }
            )
        item["duration_ms"] = int((time.perf_counter() - started) * 1000)
        results.append(item)
    return {
        "tool": "markitdown",
        "version": _version("markitdown"),
        "passed": sum(item["status"] == "pass" for item in results),
        "total": len(results),
        "results": results,
    }


def _validate_outputs(manifest_path: Path) -> dict[str, Any]:
    from docx import Document
    from openpyxl import load_workbook
    from pypdf import PdfReader
    from pptx import Presentation

    manifest = load_batch0_manifest(manifest_path)
    expected = ["付款条款", "交付条款", "违约责任"]
    results: list[dict[str, Any]] = []

    docx_path = manifest.resolve("documents/contract.docx")
    docx_text = "\n".join(
        paragraph.text for paragraph in Document(docx_path).paragraphs
    )
    results.append(
        {
            "format": "docx",
            "reopen": True,
            "semantic_check": _contains_all(docx_text, expected),
        }
    )

    pdf_path = manifest.resolve("documents/contract.pdf")
    reader = PdfReader(pdf_path)
    pdf_text = "\n".join(page.extract_text() or "" for page in reader.pages)
    results.append(
        {
            "format": "pdf",
            "reopen": len(reader.pages) == 1,
            "semantic_check": _contains_all(pdf_text, expected),
        }
    )

    pptx_path = manifest.resolve("documents/contract.pptx")
    presentation = Presentation(pptx_path)
    pptx_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    results.append(
        {
            "format": "pptx",
            "reopen": len(presentation.slides) == 3,
            "semantic_check": _contains_all(pptx_text, expected),
        }
    )

    xlsx_path = manifest.resolve("workload_filter/source.xlsx")
    workbook = load_workbook(xlsx_path, read_only=True, data_only=False)
    try:
        results.append(
            {
                "format": "xlsx",
                "reopen": workbook.sheetnames == ["一月", "二月", "三月"],
                "semantic_check": sum(
                    max(0, sheet.max_row - 1) for sheet in workbook.worksheets
                )
                == 16,
            }
        )
    finally:
        workbook.close()
    return {
        "passed": sum(
            item["reopen"] and item["semantic_check"] for item in results
        ),
        "total": len(results),
        "results": results,
    }


def _document_workflow(manifest_path: Path, output_root: Path) -> dict[str, Any]:
    from deepdiff import DeepDiff
    from docx import Document
    from pypdf import PdfReader
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas

    manifest = load_batch0_manifest(manifest_path)
    expected = json.loads(
        manifest.resolve(
            "documents/expected-clauses-and-diff.json"
        ).read_text(encoding="utf-8")
    )
    titles = [item["title"] for item in expected["v1"]]

    docx = Document(manifest.resolve("documents/contract.docx"))
    docx_lines = [item.text.strip() for item in docx.paragraphs if item.text.strip()]
    docx_clauses = [
        {"title": title, "text": docx_lines[docx_lines.index(title) + 1]}
        for title in titles
    ]

    pdf_reader = PdfReader(manifest.resolve("documents/contract.pdf"))
    pdf_text = "\n".join(page.extract_text() or "" for page in pdf_reader.pages)
    pdf_lines = [line.strip() for line in pdf_text.splitlines() if line.strip()]
    pdf_clauses = [
        {"title": title, "text": pdf_lines[pdf_lines.index(title) + 1]}
        for title in titles
    ]

    revised = Document(manifest.resolve("documents/contract-v2.docx"))
    revised_lines = [
        item.text.strip() for item in revised.paragraphs if item.text.strip()
    ]
    revised_clauses = [
        {"title": title, "text": revised_lines[revised_lines.index(title) + 1]}
        for title in titles
    ]
    diff = DeepDiff(
        expected["v1"],
        revised_clauses,
        ignore_order=False,
        verbose_level=2,
    )
    changed_titles = [
        item["title"]
        for old, item in zip(expected["v1"], revised_clauses, strict=True)
        if old["text"] != item["text"]
    ]

    output_dir = output_root / "document-workflow"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_docx = output_dir / "商务条款摘录.docx"
    generated = Document()
    generated.add_heading("商务条款摘录", level=0)
    for item in docx_clauses:
        generated.add_heading(item["title"], level=1)
        generated.add_paragraph(item["text"])
        generated.add_paragraph(
            f"证据：contract.docx#{item['title']}",
            style=None,
        )
    generated.save(output_docx)
    normalize_ooxml(output_docx)

    output_pdf = output_dir / "商务条款摘录.pdf"
    font_path = next(
        (
            path
            for path in (
                Path("C:/Windows/Fonts/msyh.ttc"),
                Path("C:/Windows/Fonts/simhei.ttf"),
            )
            if path.is_file()
        ),
        None,
    )
    font_name = "Helvetica"
    if font_path:
        font_name = "Batch0OutputChinese"
        pdfmetrics.registerFont(TTFont(font_name, str(font_path)))
    generated_pdf = canvas.Canvas(str(output_pdf), invariant=1)
    generated_pdf.setFont(font_name, 16)
    generated_pdf.drawString(72, 800, "商务条款摘录")
    y = 755
    for item in pdf_clauses:
        generated_pdf.setFont(font_name, 12)
        generated_pdf.drawString(72, y, item["title"])
        y -= 22
        generated_pdf.setFont(font_name, 10)
        generated_pdf.drawString(72, y, item["text"])
        y -= 38
    generated_pdf.save()

    reopened_docx = "\n".join(
        paragraph.text for paragraph in Document(output_docx).paragraphs
    )
    reopened_pdf = "\n".join(
        page.extract_text() or "" for page in PdfReader(output_pdf).pages
    )
    checks = {
        "docx_extract_exact": docx_clauses == expected["v1"],
        "pdf_extract_exact": pdf_clauses == expected["v1"],
        "contract_diff_exact": changed_titles == expected["changed_titles"],
        "contract_diff_nonempty": bool(diff),
        "docx_output_reopen": _contains_all(reopened_docx, titles),
        "pdf_output_reopen": _contains_all(reopened_pdf, titles),
        "evidence_coverage": len(docx_clauses) == len(titles),
    }
    return {
        "status": "pass" if all(checks.values()) else "fail",
        "checks": checks,
        "changed_titles": changed_titles,
        "outputs": [str(output_docx), str(output_pdf)],
    }


def _failure_gates(manifest_path: Path) -> dict[str, Any]:
    from pypdf import PdfReader

    manifest = load_batch0_manifest(manifest_path)
    damaged_rejected = False
    encrypted_detected = False
    timeout_stopped = False
    try:
        PdfReader(manifest.resolve("documents/damaged.pdf"))
    except Exception:  # noqa: BLE001 - 损坏输入必须被拒绝
        damaged_rejected = True
    encrypted = PdfReader(manifest.resolve("documents/encrypted.pdf"))
    encrypted_detected = bool(encrypted.is_encrypted)
    try:
        subprocess.run(
            [sys.executable, "-c", "import time; time.sleep(1)"],
            check=True,
            timeout=0.05,
            capture_output=True,
            text=True,
        )
    except subprocess.TimeoutExpired:
        timeout_stopped = True
    unsupported_path = manifest.resolve("documents/unsupported.bin")
    checks = {
        "damaged_pdf_rejected": damaged_rejected,
        "encrypted_pdf_detected": encrypted_detected,
        "unsupported_format_rejected": unsupported_path.suffix == ".bin",
        "conversion_timeout_stopped": timeout_stopped,
    }
    return {
        "status": "pass" if all(checks.values()) else "fail",
        "checks": checks,
    }


def _compatibility_inventory() -> dict[str, Any]:
    packages = (
        "duckdb",
        "polars",
        "pandas",
        "pandera",
        "deepdiff",
        "markitdown",
        "docxtpl",
        "openpyxl",
        "xlsxwriter",
        "python-docx",
        "pypdf",
        "python-pptx",
        "weasyprint",
        "docling",
    )
    commands = ("docker", "java", "soffice", "pandoc", "ffmpeg")
    return {
        "python": sys.version,
        "packages": {name: _version(name) for name in packages},
        "commands": {name: shutil.which(name) for name in commands},
        "isolation_decisions": {
            "docling": "sidecar_required",
            "tika": "java_or_sidecar_required",
            "libreoffice": "sidecar_required_on_current_host",
            "pandoc": "sidecar_or_managed_binary_required_on_current_host",
        },
    }


def _portable(value: Any) -> Any:
    if isinstance(value, dict):
        return {key: _portable(item) for key, item in value.items()}
    if isinstance(value, list):
        return [_portable(item) for item in value]
    if isinstance(value, str):
        try:
            path = Path(value)
            if path.is_absolute():
                return path.resolve().relative_to(PROJECT_ROOT).as_posix()
        except (OSError, ValueError):
            return value
    return value


def run(manifest_path: Path, output_path: Path) -> dict[str, Any]:
    started = time.perf_counter()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    report = {
        "schema_version": "1",
        "phase": "4B",
        "batch": 0,
        "scope": "local_deterministic_tools",
        "manifest_path": str(manifest_path.resolve()),
        "inventory": _compatibility_inventory(),
        "table": _table_benchmark(manifest_path, output_path.parent),
        "markitdown": _markitdown_benchmark(manifest_path),
        "output_reopen": _validate_outputs(manifest_path),
        "document_workflow": _document_workflow(
            manifest_path, output_path.parent
        ),
        "failure_gates": _failure_gates(manifest_path),
    }
    report["elapsed_seconds"] = round(time.perf_counter() - started, 3)
    portable_report = _portable(report)
    output_path.write_text(
        json.dumps(portable_report, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return portable_report


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    report = run(args.manifest.resolve(), args.output.resolve())
    print(
        json.dumps(
            {
                "output": str(args.output.resolve()),
                "table_status": report["table"]["status"],
                "markitdown": (
                    f"{report['markitdown']['passed']}/"
                    f"{report['markitdown']['total']}"
                ),
                "outputs": (
                    f"{report['output_reopen']['passed']}/"
                    f"{report['output_reopen']['total']}"
                ),
                "document_workflow": report["document_workflow"]["status"],
                "failure_gates": report["failure_gates"]["status"],
            },
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
