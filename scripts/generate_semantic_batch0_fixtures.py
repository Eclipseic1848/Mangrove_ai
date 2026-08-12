# -*- coding: utf-8 -*-
"""确定性生成 Phase 4B 批次 0 的公开脱敏 Golden。"""
from __future__ import annotations

import csv
import hashlib
import io
import json
from pathlib import Path
import re
import sys
import zipfile


PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_ROOT = (
    PROJECT_ROOT / "tests" / "fixtures" / "semantic_harness" / "public" / "batch0"
)
FIXED_ZIP_TIME = (2020, 1, 1, 0, 0, 0)


def normalize_ooxml(path: Path) -> None:
    """固定 OOXML 包内时间戳和核心属性，确保 Golden 可重复生成。"""

    source = path.read_bytes()
    output = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(source), "r") as reader:
        with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as writer:
            for old_info in sorted(reader.infolist(), key=lambda item: item.filename):
                content = reader.read(old_info.filename)
                if old_info.filename == "docProps/core.xml":
                    text = content.decode("utf-8")
                    text = re.sub(
                        r"<dcterms:(created|modified)([^>]*)>.*?</dcterms:\1>",
                        (
                            r"<dcterms:\1\2>"
                            r"2020-01-01T00:00:00Z"
                            r"</dcterms:\1>"
                        ),
                        text,
                    )
                    content = text.encode("utf-8")
                info = zipfile.ZipInfo(old_info.filename, FIXED_ZIP_TIME)
                info.compress_type = zipfile.ZIP_DEFLATED
                info.external_attr = old_info.external_attr
                info.create_system = old_info.create_system
                writer.writestr(info, content)
    path.write_bytes(output.getvalue())


def _rows() -> list[dict[str, str]]:
    names = [
        "示例人员甲",
        "示例人员乙",
        "示例人员甲",
        "示例人员甲",
        "示例人员丙",
        "示例人员甲",
        "示例人员甲",
        "示例人员甲",
        "示例人员丁",
        "示例人员甲",
        "示例人员甲",
        "示例人员甲",
        "示例人员乙",
        "示例人员甲",
        "示例人员甲",
        "示例人员戊",
    ]
    rows: list[dict[str, str]] = []
    for index, name in enumerate(names, start=1):
        rows.append(
            {
                "姓名": name,
                "核销工作量天数": f"{(index % 5 + 1) / 2:.1f}",
                "工作量费用": f"{index * 680:.2f}",
                "项目": f"示例项目{chr(64 + (index % 3) + 1)}",
                "月份": f"2026-{(index % 3) + 1:02d}",
                "__source_ref": f"source.parquet#row={index}",
            }
        )
    return rows


def _write_csv(path: Path, rows: list[dict[str, str]], delimiter: str = ",") -> None:
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0]), delimiter=delimiter)
        writer.writeheader()
        writer.writerows(rows)


def _write_table_assets(root: Path, rows: list[dict[str, str]]) -> dict[str, str]:
    import pandas as pd

    table_dir = root / "workload_filter"
    table_dir.mkdir(parents=True, exist_ok=True)
    _write_csv(table_dir / "source.csv", rows)
    _write_csv(table_dir / "source.tsv", rows, delimiter="\t")
    (table_dir / "source.json").write_text(
        json.dumps(rows, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    (table_dir / "source.jsonl").write_text(
        "\n".join(json.dumps(row, ensure_ascii=False) for row in rows) + "\n",
        encoding="utf-8",
    )
    frame = pd.DataFrame(rows)
    frame.to_parquet(table_dir / "source.parquet", index=False)
    with pd.ExcelWriter(table_dir / "source.xlsx", engine="openpyxl") as writer:
        frame.iloc[:6].to_excel(writer, sheet_name="一月", index=False)
        frame.iloc[6:11].to_excel(writer, sheet_name="二月", index=False)
        frame.iloc[11:].to_excel(writer, sheet_name="三月", index=False)
    normalize_ooxml(table_dir / "source.xlsx")
    frame.to_html(table_dir / "source.html", index=False, encoding="utf-8")

    expected = [
        {
            "核销工作量天数": row["核销工作量天数"],
            "工作量费用": row["工作量费用"],
        }
        for row in rows
        if row["姓名"] == "示例人员甲"
    ]
    (table_dir / "expected-records.json").write_text(
        json.dumps(expected, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    _write_csv(table_dir / "expected.csv", expected)
    return {
        "csv": "workload_filter/source.csv",
        "tsv": "workload_filter/source.tsv",
        "xlsx": "workload_filter/source.xlsx",
        "json": "workload_filter/source.json",
        "jsonl": "workload_filter/source.jsonl",
        "parquet": "workload_filter/source.parquet",
        "html": "workload_filter/source.html",
    }


def _write_table_operation_assets(root: Path) -> None:
    import pandas as pd

    operation_dir = root / "table_operations"
    operation_dir.mkdir(parents=True, exist_ok=True)
    facts = [
        {
            "记录号": "R1",
            "姓名": "示例人员甲",
            "部门": "研发一部",
            "费用": "100.00",
            "__source_ref": "facts.parquet#row=1",
        },
        {
            "记录号": "R2",
            "姓名": "示例人员乙",
            "部门": "研发一部",
            "费用": "200.00",
            "__source_ref": "facts.parquet#row=2",
        },
        {
            "记录号": "R3",
            "姓名": "示例人员丙",
            "部门": "交付一部",
            "费用": "300.00",
            "__source_ref": "facts.parquet#row=3",
        },
        {
            "记录号": "R4",
            "姓名": "示例人员丁",
            "部门": "交付一部",
            "费用": "400.00",
            "__source_ref": "facts.parquet#row=4",
        },
    ]
    departments = [
        {"部门": "研发一部", "区域": "华东"},
        {"部门": "交付一部", "区域": "华南"},
    ]
    duplicated_departments = [
        *departments,
        {"部门": "研发一部", "区域": "华北"},
    ]
    pd.DataFrame(facts).to_parquet(operation_dir / "facts.parquet", index=False)
    pd.DataFrame(departments).to_parquet(
        operation_dir / "departments.parquet", index=False
    )
    pd.DataFrame(duplicated_departments).to_parquet(
        operation_dir / "departments-duplicate.parquet", index=False
    )
    expected = {
        "merge": [
            {"记录号": "R1", "姓名": "示例人员甲", "区域": "华东"},
            {"记录号": "R2", "姓名": "示例人员乙", "区域": "华东"},
            {"记录号": "R3", "姓名": "示例人员丙", "区域": "华南"},
            {"记录号": "R4", "姓名": "示例人员丁", "区域": "华南"},
        ],
        "aggregate": [
            {"部门": "交付一部", "费用合计": "700.00"},
            {"部门": "研发一部", "费用合计": "300.00"},
        ],
        "negative_cases": [
            {
                "case_id": "missing_projection_column",
                "expected_error": "missing_column",
            },
            {
                "case_id": "non_unique_join_key",
                "expected_error": "join_cardinality_violation",
            },
        ],
    }
    (operation_dir / "expected.json").write_text(
        json.dumps(expected, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def _find_chinese_font() -> Path | None:
    candidates = (
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/simsun.ttc"),
    )
    return next((path for path in candidates if path.is_file()), None)


def _write_document_assets(root: Path) -> None:
    from docx import Document
    from pptx import Presentation
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas

    document_dir = root / "documents"
    document_dir.mkdir(parents=True, exist_ok=True)
    clauses = [
        ("付款条款", "验收通过后十五个工作日内支付合同金额的百分之六十。"),
        ("交付条款", "乙方应在二零二六年九月三十日前完成全部成果交付。"),
        ("违约责任", "逾期交付的，每日按未交付部分金额的千分之一承担违约责任。"),
    ]
    revised_clauses = [
        ("付款条款", "验收通过后十五个工作日内支付合同金额的百分之七十。"),
        ("交付条款", "乙方应在二零二六年十月十五日前完成全部成果交付。"),
        ("违约责任", "逾期交付的，每日按未交付部分金额的千分之二承担违约责任。"),
    ]
    markdown = "# 示例商务合同\n\n" + "\n\n".join(
        f"## {title}\n\n{text}" for title, text in clauses
    )
    (document_dir / "contract.md").write_text(markdown + "\n", encoding="utf-8")
    (document_dir / "contract.txt").write_text(
        "\n".join(f"{title}：{text}" for title, text in clauses) + "\n",
        encoding="utf-8",
    )
    (document_dir / "contract.html").write_text(
        "<!doctype html><html lang=\"zh-CN\"><meta charset=\"utf-8\">"
        "<title>示例商务合同</title><body><h1>示例商务合同</h1>"
        + "".join(f"<h2>{title}</h2><p>{text}</p>" for title, text in clauses)
        + "</body></html>",
        encoding="utf-8",
    )
    (document_dir / "contract.xml").write_text(
        "<?xml version=\"1.0\" encoding=\"UTF-8\"?>\n<contract>"
        + "".join(
            f"<clause title=\"{title}\">{text}</clause>" for title, text in clauses
        )
        + "</contract>\n",
        encoding="utf-8",
    )

    document = Document()
    document.add_heading("示例商务合同", level=0)
    for title, text in clauses:
        document.add_heading(title, level=1)
        document.add_paragraph(text)
    document.save(document_dir / "contract.docx")
    normalize_ooxml(document_dir / "contract.docx")
    revised_document = Document()
    revised_document.add_heading("示例商务合同（修订版）", level=0)
    for title, text in revised_clauses:
        revised_document.add_heading(title, level=1)
        revised_document.add_paragraph(text)
    revised_document.save(document_dir / "contract-v2.docx")
    normalize_ooxml(document_dir / "contract-v2.docx")

    font_path = _find_chinese_font()
    font_name = "Helvetica"
    if font_path:
        font_name = "Batch0Chinese"
        pdfmetrics.registerFont(TTFont(font_name, str(font_path)))
    pdf = canvas.Canvas(str(document_dir / "contract.pdf"), invariant=1)
    pdf.setTitle("示例商务合同")
    y = 800
    pdf.setFont(font_name, 18)
    pdf.drawString(72, y, "示例商务合同")
    for title, text in clauses:
        y -= 50
        pdf.setFont(font_name, 14)
        pdf.drawString(72, y, title)
        y -= 24
        pdf.setFont(font_name, 11)
        pdf.drawString(72, y, text)
    pdf.save()
    revised_pdf = canvas.Canvas(str(document_dir / "contract-v2.pdf"), invariant=1)
    revised_pdf.setTitle("示例商务合同（修订版）")
    y = 800
    revised_pdf.setFont(font_name, 18)
    revised_pdf.drawString(72, y, "示例商务合同（修订版）")
    for title, text in revised_clauses:
        y -= 50
        revised_pdf.setFont(font_name, 14)
        revised_pdf.drawString(72, y, title)
        y -= 24
        revised_pdf.setFont(font_name, 11)
        revised_pdf.drawString(72, y, text)
    revised_pdf.save()

    presentation = Presentation()
    for title, text in clauses:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = text
    presentation.save(document_dir / "contract.pptx")
    normalize_ooxml(document_dir / "contract.pptx")

    damaged = (document_dir / "contract.pdf").read_bytes()[:128]
    (document_dir / "damaged.pdf").write_bytes(damaged)
    (document_dir / "unsupported.bin").write_bytes(b"\x00BATCH0-UNSUPPORTED\xff")
    clause_payload = {
        "v1": [{"title": title, "text": text} for title, text in clauses],
        "v2": [{"title": title, "text": text} for title, text in revised_clauses],
        "changed_titles": [title for title, _ in clauses],
    }
    (document_dir / "expected-clauses-and-diff.json").write_text(
        json.dumps(clause_payload, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )

    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(document_dir / "contract.pdf")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt("batch0")
    with (document_dir / "encrypted.pdf").open("wb") as handle:
        writer.write(handle)


def _write_image_and_zip(root: Path) -> None:
    from PIL import Image, ImageDraw, ImageFont

    assets = root / "assets"
    assets.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (1200, 500), "white")
    draw = ImageDraw.Draw(image)
    font_path = _find_chinese_font()
    font = ImageFont.truetype(str(font_path), 42) if font_path else None
    draw.text(
        (60, 70),
        "示例表单\n姓名：示例人员甲\n核销工作量天数：1.5\n工作量费用：2040.00",
        fill="black",
        font=font,
        spacing=18,
    )
    for extension in ("png", "jpg", "webp", "tiff", "bmp"):
        image.save(assets / f"form.{extension}")

    with zipfile.ZipFile(
        assets / "batch0-safe.zip", "w", zipfile.ZIP_DEFLATED
    ) as archive:
        for source, name in (
            (root / "documents" / "contract.docx", "contract.docx"),
            (root / "workload_filter" / "source.xlsx", "source.xlsx"),
        ):
            info = zipfile.ZipInfo(name, FIXED_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            archive.writestr(info, source.read_bytes())


def _file_entry(root: Path, path: Path, role: str) -> dict[str, object]:
    raw = path.read_bytes()
    return {
        "path": path.relative_to(root).as_posix(),
        "format": path.suffix.lstrip(".").lower(),
        "role": role,
        "sha256": hashlib.sha256(raw).hexdigest(),
        "size_bytes": len(raw),
    }


def main() -> int:
    OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)
    rows = _rows()
    input_paths = _write_table_assets(OUTPUT_ROOT, rows)
    _write_table_operation_assets(OUTPUT_ROOT)
    _write_document_assets(OUTPUT_ROOT)
    _write_image_and_zip(OUTPUT_ROOT)

    files = []
    for path in sorted(OUTPUT_ROOT.rglob("*")):
        if path.is_file() and path.name != "manifest.json":
            role = "expected" if "expected" in path.name else "input"
            files.append(_file_entry(OUTPUT_ROOT, path, role))
    manifest = {
        "manifest_version": "1",
        "generated_by": "scripts/generate_semantic_batch0_fixtures.py",
        "deidentified": True,
        "core_formats": [
            "xlsx",
            "csv",
            "tsv",
            "json",
            "jsonl",
            "parquet",
            "pdf",
            "docx",
            "pptx",
            "html",
            "markdown",
        ],
        "compatibility_formats": [
            "xls",
            "ods",
            "doc",
            "odt",
            "rtf",
            "txt",
            "xml",
            "epub",
            "ppt",
            "odp",
            "png",
            "jpg",
            "webp",
            "tiff",
            "bmp",
            "zip",
            "rar",
            "7z",
        ],
        "files": files,
        "cases": [
            {
                "case_id": "workload_filter",
                "task_family": "tabular_transform",
                "plan_id": "plan_workload_filter_v1",
                "canonical_input": "workload_filter/source.parquet",
                "input_paths": input_paths,
                "selection": {"姓名": "示例人员甲"},
                "projection": ["核销工作量天数", "工作量费用"],
                "expected": {
                    "records_path": "workload_filter/expected-records.json",
                    "row_count": 11,
                    "visible_columns": ["核销工作量天数", "工作量费用"],
                    "table_count": 1,
                    "evidence_coverage": 1.0,
                },
            }
        ],
    }
    (OUTPUT_ROOT / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    print(f"已生成批次 0 Golden：{OUTPUT_ROOT}")
    print(f"文件数：{len(files)}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
