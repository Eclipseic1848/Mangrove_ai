# -*- coding: utf-8 -*-
"""PPTX 只读解析器；复用 python-pptx 保留幻灯片和 Shape 结构位置。"""
from __future__ import annotations

import hashlib
import io
from importlib.metadata import PackageNotFoundError, version
import json
from typing import Any, Dict, List, Tuple

from src.data_prep.document_models import DocumentElement, ElementType
from src.data_prep.models import RawArtifact, RecordEnvelope

from .registry import Parser


try:
    _PPTX_VERSION = version("python-pptx")
except PackageNotFoundError:  # pragma: no cover
    _PPTX_VERSION = "unknown"


def _record(
    artifact: RawArtifact,
    *,
    data: Dict[str, Any],
    parser_name: str,
    position: Dict[str, Any],
) -> RecordEnvelope:
    content_hash = hashlib.sha256(
        json.dumps(data, ensure_ascii=False, sort_keys=True, default=str).encode("utf-8")
    ).hexdigest()
    record_id = hashlib.sha256(
        f"{artifact.artifact_id}:{parser_name}:{position}".encode("utf-8")
    ).hexdigest()[:16]
    return RecordEnvelope(
        record_id=record_id,
        data=data,
        meta={
            "source_id": artifact.source_id,
            "artifact_id": artifact.artifact_id,
            "parser": parser_name,
            "position": position,
            "content_hash": content_hash,
        },
    )


def _element(
    artifact: RawArtifact,
    *,
    page: int,
    element_type: ElementType,
    text: str,
    reading_order: int,
    location: Dict[str, Any],
    metadata: Dict[str, Any] | None = None,
) -> Dict[str, Any]:
    seed = json.dumps(
        {
            "artifact_id": artifact.artifact_id,
            "type": element_type.value,
            "location": location,
            "text": text,
        },
        ensure_ascii=False,
        sort_keys=True,
    )
    return DocumentElement(
        element_id=f"el-{hashlib.sha256(seed.encode('utf-8')).hexdigest()[:20]}",
        artifact_id=artifact.artifact_id,
        page=page,
        element_type=element_type,
        text=text,
        reading_order=reading_order,
        extractor="python-pptx",
        extractor_version=_PPTX_VERSION,
        metadata={"location": location, **(metadata or {})},
    ).model_dump(mode="json")


class PresentationParser(Parser):
    """提取 PPTX 文本框、占位符和表格，不执行宏或外部对象。"""

    name = "pptx"
    media_types = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    extensions = ("pptx",)

    def parse(
        self, artifact: RawArtifact, raw_bytes: bytes
    ) -> Tuple[List[RecordEnvelope], List[Dict]]:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER

        try:
            presentation = Presentation(io.BytesIO(raw_bytes))
        except Exception as exc:  # noqa: BLE001
            return [], [{
                "artifact_id": artifact.artifact_id,
                "reason": f"pptx_parse_failed: {exc}",
                "position": {},
            }]

        records: List[RecordEnvelope] = []
        reading_order = 0
        for slide_no, slide in enumerate(presentation.slides, start=1):
            shapes = sorted(
                slide.shapes,
                key=lambda shape: (
                    int(getattr(shape, "top", 0)),
                    int(getattr(shape, "left", 0)),
                    int(getattr(shape, "shape_id", 0)),
                ),
            )
            for shape in shapes:
                shape_id = int(getattr(shape, "shape_id", 0))
                if getattr(shape, "has_table", False):
                    table = shape.table
                    header = [cell.text.strip() for cell in table.rows[0].cells]
                    for row_no, row in enumerate(table.rows[1:], start=2):
                        cells = [cell.text.strip() for cell in row.cells]
                        if not any(cells):
                            continue
                        data = {
                            (header[index] if index < len(header) and header[index] else f"col_{index}"): (
                                cells[index] if index < len(cells) else ""
                            )
                            for index in range(max(len(header), len(cells)))
                        }
                        text = "；".join(
                            f"{key}：{value}" for key, value in data.items()
                        )
                        location = {
                            "kind": "pptx_table_row",
                            "slide": slide_no,
                            "shape_id": shape_id,
                            "row": row_no,
                        }
                        data["elements"] = [_element(
                            artifact,
                            page=slide_no,
                            element_type=ElementType.TABLE,
                            text=text,
                            reading_order=reading_order,
                            location=location,
                            metadata={
                                "table_columns": list(data),
                                "table_row": dict(data),
                            },
                        )]
                        records.append(_record(
                            artifact,
                            data=data,
                            parser_name="pptx_table",
                            position=location,
                        ))
                        reading_order += 1
                    continue

                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph_no, paragraph in enumerate(
                    shape.text_frame.paragraphs,
                    start=1,
                ):
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    location = {
                        "kind": "pptx_paragraph",
                        "slide": slide_no,
                        "shape_id": shape_id,
                        "paragraph": paragraph_no,
                    }
                    records.append(_record(
                        artifact,
                        data={
                            "text": text,
                            "elements": [_element(
                                artifact,
                                page=slide_no,
                                element_type=(
                                    ElementType.HEADING
                                    if getattr(shape, "is_placeholder", False)
                                    and shape.placeholder_format.type in {
                                        PP_PLACEHOLDER.TITLE,
                                        PP_PLACEHOLDER.CENTER_TITLE,
                                        PP_PLACEHOLDER.SUBTITLE,
                                    }
                                    and paragraph_no == 1
                                    else ElementType.PARAGRAPH
                                ),
                                text=text,
                                reading_order=reading_order,
                                location=location,
                                metadata={
                                    "level": int(getattr(paragraph, "level", 0)),
                                    "shape_name": getattr(shape, "name", ""),
                                },
                            )],
                        },
                        parser_name="pptx",
                        position=location,
                    ))
                    reading_order += 1
        return records, []
