# -*- coding: utf-8 -*-
"""基于成熟文件库的确定性 Renderer、QA 与原子发布。"""
from __future__ import annotations

import csv
from dataclasses import dataclass
import hashlib
import html
from html.parser import HTMLParser
from importlib.metadata import version
import json
import mimetypes
from pathlib import Path
import re
import shutil
from typing import Any, Callable, Mapping, Sequence
import uuid

import pyarrow as pa
import pyarrow.parquet as pq

from ..document_models import DocumentExecutionResult
from ..models import DeliveryFormat, SemanticTaskPlan
from .models import (
    ArtifactQAReport,
    DeliveryManifest,
    DeliveryOutput,
    DeliveryStatus,
)


_FORMATS = (
    DeliveryFormat.JSON,
    DeliveryFormat.JSONL,
    DeliveryFormat.CSV,
    DeliveryFormat.XLSX,
    DeliveryFormat.PARQUET,
    DeliveryFormat.DOCX,
    DeliveryFormat.PDF,
    DeliveryFormat.HTML,
    DeliveryFormat.MARKDOWN,
    DeliveryFormat.TXT,
    DeliveryFormat.PPTX,
)
_EXTENSIONS = {
    DeliveryFormat.MARKDOWN: "md",
}
_MEDIA_TYPES = {
    DeliveryFormat.JSON: "application/json",
    DeliveryFormat.JSONL: "application/x-ndjson",
    DeliveryFormat.CSV: "text/csv; charset=utf-8",
    DeliveryFormat.XLSX: (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ),
    DeliveryFormat.PARQUET: "application/vnd.apache.parquet",
    DeliveryFormat.DOCX: (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ),
    DeliveryFormat.PDF: "application/pdf",
    DeliveryFormat.HTML: "text/html; charset=utf-8",
    DeliveryFormat.MARKDOWN: "text/markdown; charset=utf-8",
    DeliveryFormat.TXT: "text/plain; charset=utf-8",
    DeliveryFormat.PPTX: (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ),
}


@dataclass(frozen=True)
class CanonicalContent:
    title: str
    columns: tuple[str, ...]
    rows: tuple[dict[str, Any], ...]
    sections: tuple[tuple[str, str], ...]
    raw: dict[str, Any]


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _display(value: Any) -> Any:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    return json.dumps(value, ensure_ascii=False, sort_keys=True, default=str)


def _safe_name(value: str | None) -> str:
    value = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "-", value or "交付结果")
    value = value.strip(" .-")
    return value[:80] or "交付结果"


def _passage_sections(passages: Sequence[Any]) -> tuple[tuple[str, str], ...]:
    """同一语义目标只展示一次标题，原文仍按来源顺序完整保留。"""

    grouped: dict[str, list[str]] = {}
    for item in passages:
        grouped.setdefault(item.label, []).append(item.text)
    return tuple(
        (label, "\n\n".join(texts))
        for label, texts in grouped.items()
    )


def _content(
    plan: SemanticTaskPlan,
    artifact_paths: Mapping[str, Path],
) -> CanonicalContent:
    parquet_path = artifact_paths.get("result")
    if parquet_path and parquet_path.suffix.lower() == ".parquet":
        table = pq.read_table(parquet_path)
        public_columns = [
            name
            for name in table.column_names
            if not name.startswith("__mg_")
        ]
        table = table.select(public_columns)
        rows = tuple(table.to_pylist())
        columns = tuple(table.column_names)
        return CanonicalContent(
            title=plan.delivery.output_name or "表格交付结果",
            columns=columns,
            rows=rows,
            sections=(),
            raw={"columns": list(columns), "rows": rows},
        )
    document_path = artifact_paths.get("result")
    if document_path is None:
        raise ValueError("缺少已验证的权威结果路径")
    result = DocumentExecutionResult.model_validate_json(
        document_path.read_text(encoding="utf-8")
    )
    sections: list[tuple[str, str]] = []
    sections.extend(_passage_sections(result.passages))
    sections.extend(
        (
            item.label,
            f"{item.change_type}: {item.before or ''} → {item.after or ''}",
        )
        for item in result.differences
    )
    sections.extend((item.label, item.message) for item in result.findings)
    sections.extend(
        (item.action.value, item.content) for item in result.derived_content
    )
    if not sections:
        sections.append(("结果", "文档执行已完成，详见结构化 JSON。"))
    raw = result.model_dump(mode="json")
    rows = tuple(
        {"标题": title, "内容": text} for title, text in sections
    )
    return CanonicalContent(
        title=plan.delivery.output_name or "文档交付结果",
        columns=("标题", "内容"),
        rows=rows,
        sections=tuple(sections),
        raw=raw,
    )


def _text(content: CanonicalContent) -> str:
    if content.sections:
        return "\n\n".join(
            f"{title}\n{text}" for title, text in content.sections
        )
    lines = [content.title, "\t".join(content.columns)]
    lines.extend(
        "\t".join(
            "" if row.get(column) is None else str(row.get(column))
            for column in content.columns
        )
        for row in content.rows
    )
    return "\n".join(lines)


def _markdown(content: CanonicalContent) -> str:
    if content.sections:
        blocks = [f"# {content.title}"]
        blocks.extend(f"## {title}\n\n{text}" for title, text in content.sections)
        return "\n\n".join(blocks) + "\n"
    header = "| " + " | ".join(content.columns) + " |"
    separator = "| " + " | ".join("---" for _ in content.columns) + " |"
    body = [
        "| "
        + " | ".join(
            str(row.get(column, "")).replace("|", "\\|").replace("\n", " ")
            for column in content.columns
        )
        + " |"
        for row in content.rows
    ]
    return "\n".join([f"# {content.title}", "", header, separator, *body]) + "\n"


def _render_json(path: Path, content: CanonicalContent) -> None:
    path.write_text(
        json.dumps(content.raw, ensure_ascii=False, indent=2, default=str),
        encoding="utf-8",
    )


def _render_jsonl(path: Path, content: CanonicalContent) -> None:
    with path.open("w", encoding="utf-8", newline="\n") as handle:
        for row in content.rows:
            handle.write(json.dumps(row, ensure_ascii=False, default=str) + "\n")


def _render_csv(path: Path, content: CanonicalContent) -> None:
    with path.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=content.columns)
        writer.writeheader()
        writer.writerows(content.rows)


def _render_parquet(path: Path, content: CanonicalContent) -> None:
    pq.write_table(pa.Table.from_pylist(list(content.rows)), path)


def _render_xlsx(path: Path, content: CanonicalContent) -> None:
    import xlsxwriter

    workbook = xlsxwriter.Workbook(
        path,
        {
            "constant_memory": True,
            "strings_to_formulas": False,
            "strings_to_urls": False,
        },
    )
    sheet = workbook.add_worksheet("结果")
    header_format = workbook.add_format({"bold": True, "bg_color": "#D9EAF7"})
    for index, column in enumerate(content.columns):
        sheet.write(0, index, column, header_format)
    for row_number, row in enumerate(content.rows, start=1):
        for column_number, column in enumerate(content.columns):
            value = _display(row.get(column))
            sheet.write(row_number, column_number, "" if value is None else value)
    sheet.freeze_panes(1, 0)
    sheet.autofilter(0, 0, max(0, len(content.rows)), len(content.columns) - 1)
    workbook.close()


def _render_docx(path: Path, content: CanonicalContent) -> None:
    from docx import Document

    document = Document()
    document.add_heading(content.title, level=0)
    if content.sections:
        for title, text in content.sections:
            document.add_heading(title, level=1)
            document.add_paragraph(text)
    else:
        table = document.add_table(rows=1, cols=len(content.columns))
        table.style = "Table Grid"
        for index, column in enumerate(content.columns):
            table.rows[0].cells[index].text = column
        for row in content.rows:
            cells = table.add_row().cells
            for index, column in enumerate(content.columns):
                cells[index].text = str(row.get(column, ""))
    document.save(path)


def _render_pdf(path: Path, content: CanonicalContent) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
    from reportlab.lib.styles import getSampleStyleSheet

    font = "STSong-Light"
    pdfmetrics.registerFont(UnicodeCIDFont(font))
    for candidate in (
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
    ):
        if candidate.exists():
            pdfmetrics.registerFont(TTFont("MangroveCJK", str(candidate)))
            font = "MangroveCJK"
            break
    styles = getSampleStyleSheet()
    for style in styles.byName.values():
        style.fontName = font
    story = [Paragraph(html.escape(content.title), styles["Title"]), Spacer(1, 12)]
    for title, text in (
        content.sections or (("数据", _text(content)),)
    ):
        story.extend(
            [
                Paragraph(html.escape(title), styles["Heading2"]),
                Paragraph(html.escape(text).replace("\n", "<br/>"), styles["BodyText"]),
                Spacer(1, 8),
            ]
        )
    SimpleDocTemplate(str(path), pagesize=A4).build(story)


def _render_html(path: Path, content: CanonicalContent) -> None:
    body = [
        "<!doctype html><html lang=\"zh-CN\"><meta charset=\"utf-8\">",
        f"<title>{html.escape(content.title)}</title>",
        "<style>body{font-family:system-ui;margin:40px;line-height:1.6}"
        "table{border-collapse:collapse}th,td{border:1px solid #ccc;padding:6px}"
        "th{background:#eef}</style>",
        f"<h1>{html.escape(content.title)}</h1>",
    ]
    if content.sections:
        body.extend(
            f"<section><h2>{html.escape(title)}</h2>"
            f"<p>{html.escape(text).replace(chr(10), '<br>')}</p></section>"
            for title, text in content.sections
        )
    else:
        body.append("<table><thead><tr>")
        body.extend(f"<th>{html.escape(col)}</th>" for col in content.columns)
        body.append("</tr></thead><tbody>")
        for row in content.rows:
            body.append("<tr>")
            body.extend(
                f"<td>{html.escape(str(row.get(col, '')))}</td>"
                for col in content.columns
            )
            body.append("</tr>")
        body.append("</tbody></table>")
    body.append("</html>")
    path.write_text("".join(body), encoding="utf-8")


def _render_pptx(path: Path, content: CanonicalContent) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    chunks = content.sections or tuple(
        (f"数据 {index + 1}", json.dumps(row, ensure_ascii=False, default=str))
        for index, row in enumerate(content.rows)
    )
    for title, text in chunks or (("结果", "无数据"),):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title[:100]
        frame = slide.placeholders[1].text_frame
        frame.text = text[:4000]
        for paragraph in frame.paragraphs:
            paragraph.font.size = Pt(18)
    presentation.save(path)


_RENDERERS: dict[DeliveryFormat, Callable[[Path, CanonicalContent], None]] = {
    DeliveryFormat.JSON: _render_json,
    DeliveryFormat.JSONL: _render_jsonl,
    DeliveryFormat.CSV: _render_csv,
    DeliveryFormat.XLSX: _render_xlsx,
    DeliveryFormat.PARQUET: _render_parquet,
    DeliveryFormat.DOCX: _render_docx,
    DeliveryFormat.PDF: _render_pdf,
    DeliveryFormat.HTML: _render_html,
    DeliveryFormat.MARKDOWN: lambda path, value: path.write_text(
        _markdown(value), encoding="utf-8"
    ),
    DeliveryFormat.TXT: lambda path, value: path.write_text(
        _text(value), encoding="utf-8"
    ),
    DeliveryFormat.PPTX: _render_pptx,
}


def _qa(path: Path, fmt: DeliveryFormat) -> ArtifactQAReport:
    row_count = sheet_count = page_count = slide_count = None
    checks = ["non_empty", "sha256"]
    if fmt == DeliveryFormat.JSON:
        json.loads(path.read_text(encoding="utf-8"))
    elif fmt == DeliveryFormat.JSONL:
        row_count = sum(
            1 for line in path.read_text(encoding="utf-8").splitlines()
            if line and json.loads(line) is not None
        )
    elif fmt == DeliveryFormat.CSV:
        with path.open(encoding="utf-8-sig", newline="") as handle:
            row_count = max(0, sum(1 for _ in csv.reader(handle)) - 1)
    elif fmt == DeliveryFormat.PARQUET:
        row_count = pq.read_metadata(path).num_rows
    elif fmt == DeliveryFormat.XLSX:
        from openpyxl import load_workbook
        workbook = load_workbook(path, read_only=True, data_only=True)
        try:
            sheet_count = len(workbook.worksheets)
            row_count = max(0, workbook.worksheets[0].max_row - 1)
        finally:
            workbook.close()
    elif fmt == DeliveryFormat.DOCX:
        from docx import Document
        Document(path)
    elif fmt == DeliveryFormat.PDF:
        from pypdf import PdfReader
        page_count = len(PdfReader(path).pages)
    elif fmt == DeliveryFormat.PPTX:
        from pptx import Presentation
        slide_count = len(Presentation(path).slides)
    elif fmt == DeliveryFormat.HTML:
        parser = HTMLParser()
        parser.feed(path.read_text(encoding="utf-8"))
        parser.close()
    else:
        text = path.read_text(encoding="utf-8")
        if not text.strip():
            raise ValueError(f"{fmt.value} 输出为空")
    checks.append("reopened")
    return ArtifactQAReport(
        format=fmt,
        openable=True,
        sha256=_sha256(path),
        size_bytes=path.stat().st_size,
        row_count=row_count,
        sheet_count=sheet_count,
        page_count=page_count,
        slide_count=slide_count,
        checks=tuple(checks),
    )


def qa_delivery_artifact(
    path: Path,
    fmt: DeliveryFormat,
) -> ArtifactQAReport:
    """供统一 Publisher 使用的独立重开 QA；不信任候选自报结果。"""

    return _qa(path, fmt)


def _renderer_versions() -> dict[str, str]:
    packages = {
        "json": "stdlib",
        "jsonl": "stdlib",
        "csv": "stdlib",
        "xlsx": "XlsxWriter",
        "parquet": "pyarrow",
        "docx": "python-docx",
        "pdf": "reportlab",
        "html": "stdlib",
        "markdown": "stdlib",
        "txt": "stdlib",
        "pptx": "python-pptx",
    }
    resolved: dict[str, str] = {}
    for fmt, package in packages.items():
        resolved[fmt] = (
            package
            if package == "stdlib"
            else f"{package} {version(package)}"
        )
    return resolved


def create_delivery(
    *,
    store: Any,
    output_root: Path,
    user_id: str,
    run_id: str,
    plan: SemanticTaskPlan,
    artifact_paths: Mapping[str, Path],
    source_artifact_hashes: Mapping[str, str] | None = None,
) -> DeliveryManifest:
    """全部格式通过独立 QA 后才把 staging 原子发布为正式目录。"""

    requested = tuple(dict.fromkeys(plan.delivery.formats))
    if (
        plan.delivery.requested_file_count is not None
        and plan.delivery.requested_file_count != len(requested)
    ):
        raise ValueError("请求文件数必须与正式输出格式数一致")
    unsupported = set(requested) - set(_FORMATS)
    if unsupported:
        names = ", ".join(sorted(item.value for item in unsupported))
        raise ValueError(f"批次 6 正式交付不支持：{names}")
    if hasattr(store, "latest_semantic_delivery"):
        existing = store.latest_semantic_delivery(user_id, run_id)
        if existing is not None:
            return DeliveryManifest.model_validate(
                {**existing, "user_id": user_id}
            )
    delivery_id = f"delivery_{uuid.uuid4().hex[:16]}"
    safe_user = hashlib.sha256(user_id.encode("utf-8")).hexdigest()[:16]
    base = output_root / safe_user / plan.plan_id / run_id / "delivery"
    staging = base / f".{delivery_id}.staging"
    final_dir = base / delivery_id
    staging.mkdir(parents=True, exist_ok=False)
    content = _content(plan, artifact_paths)
    outputs: list[DeliveryOutput] = []
    try:
        for fmt in requested:
            output_id = f"output_{uuid.uuid4().hex[:16]}"
            extension = _EXTENSIONS.get(fmt, fmt.value)
            filename = f"{_safe_name(plan.delivery.output_name)}.{extension}"
            path = staging / filename
            _RENDERERS[fmt](path, content)
            report = _qa(path, fmt)
            outputs.append(
                DeliveryOutput(
                    output_id=output_id,
                    format=fmt,
                    filename=filename,
                    media_type=_MEDIA_TYPES.get(
                        fmt,
                        mimetypes.guess_type(filename)[0]
                        or "application/octet-stream",
                    ),
                    sha256=report.sha256,
                    size_bytes=report.size_bytes,
                    qa=report,
                    download_url=f"/api/semantic-deliveries/outputs/{output_id}",
                )
            )
        manifest = DeliveryManifest(
            delivery_id=delivery_id,
            run_id=run_id,
            plan_id=plan.plan_id,
            user_id=user_id,
            status=DeliveryStatus.SUCCEEDED,
            source_artifact_hashes=dict(source_artifact_hashes or {}),
            requested_formats=requested,
            outputs=tuple(outputs),
            renderer_versions=_renderer_versions(),
            provenance={
                "authoritative_result_sha256": _sha256(
                    artifact_paths["result"]
                ),
                "authoritative_result_kind": (
                    "table"
                    if artifact_paths["result"].suffix.lower() == ".parquet"
                    else "document"
                ),
            },
        )
        manifest_path = staging / "manifest.json"
        manifest_path.write_text(
            manifest.model_dump_json(
                indent=2, exclude={"user_id"}
            ),
            encoding="utf-8",
        )
        final_dir.parent.mkdir(parents=True, exist_ok=True)
        staging.replace(final_dir)
        store.save_semantic_delivery(
            user_id=user_id,
            run_id=run_id,
            manifest=manifest,
            output_dir=final_dir,
        )
        return manifest
    except Exception:
        if staging.exists():
            shutil.rmtree(staging)
        if final_dir.exists():
            shutil.rmtree(final_dir)
        if hasattr(store, "latest_semantic_delivery"):
            existing = store.latest_semantic_delivery(user_id, run_id)
            if existing is not None:
                return DeliveryManifest.model_validate(
                    {**existing, "user_id": user_id}
                )
        raise
