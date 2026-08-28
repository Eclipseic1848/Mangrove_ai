# -*- coding: utf-8 -*-
"""Phase 4B 批次 7：正式数据工作台 API。"""
from __future__ import annotations

import asyncio
from dataclasses import dataclass
from datetime import datetime, timezone
import hashlib
import json
from pathlib import Path
import re
import sqlite3
import uuid
import zipfile
from typing import Any, Callable, Literal

import duckdb
from fastapi import (
    APIRouter,
    Depends,
    Header,
    HTTPException,
    Query,
    status,
)
from fastapi.responses import FileResponse
from pydantic import (
    BaseModel,
    ConfigDict,
    Field,
    field_validator,
    model_validator,
)
from sse_starlette.sse import EventSourceResponse
from starlette.background import BackgroundTask

from src.agentic_runtime.models import (
    PermissionProfile,
    RuntimeTaskConfig,
    RuntimeVersion,
)
from src.agentic_runtime.repository import AgenticRuntimeRepository
from src.api.auth import get_current_user, get_store, is_admin_role
from src.api.catalog_actor import catalog_actor_from_user
from src.api.semantic_workspace_runtime import (
    get_semantic_workspace_manager,
)
from src.capability_catalog import (
    CapabilityCatalog,
    CapabilityMountGateRejected,
    CapabilityPackRef,
    DefaultCapabilityMounts,
    SqliteCapabilityCatalogRepository,
)
from src.config.settings import settings
from src.candidate_verification import (
    HistoricalAuthorityRecoveryConfirmation,
    ReverificationContractError,
    ReverificationUnavailableError,
    RulesetIdentityStatus,
    SqliteCandidateVerificationRepository,
)
from src.conversation_steering import (
    CapabilityMaturity,
    ConversationSteering,
    ProgressAudience,
    ProgressProjection,
    ProgressStage,
    ProgressValue,
    RevisionDecisionStatus,
    RevisionProposalStatus,
    RevisionSwitchMode,
    SqliteSteeringRepository,
    SteeringAction,
    SteeringRequest,
    StructuredProgressEvent,
    build_context_rewriter,
)
from src.model_connections import GrantError, get_default_broker
from src.delivery_publishing.models import TableOutputContract
from src.services.upload_store import UploadStore
from src.source_acquisition import (
    AcquisitionConflictError,
    AnonymousWebFetcher,
    SourceAcquisitionRepository,
    SourceAcquisitionRequest,
    SourceAcquisitionService,
)
from src.runtime_routing import (
    RolloutActor,
    RolloutSnapshot,
    RuntimeAssignment,
    RuntimeRouting,
    RuntimeTaskRevisionRef,
    SqliteRuntimeRoutingRepository,
    open_runtime_routing_repository,
)
from src.task_context import (
    TaskContextRepository,
    TaskContextSelection,
    TaskContextService,
)


router = APIRouter(
    prefix="/api/semantic-workspace",
    tags=["semantic-workspace"],
)

_FORMATS = {
    "json",
    "jsonl",
    "csv",
    "xlsx",
    "parquet",
    "docx",
    "pdf",
    "html",
    "markdown",
    "txt",
    "pptx",
}
_DOCUMENT_INPUTS = {".docx", ".pdf"}
_TABLE_INPUTS = {".xlsx", ".csv", ".tsv", ".json", ".jsonl", ".parquet"}
_OUTPUT_FORMAT_PATTERN = re.compile(
    r"(?:输出|导出|生成)(?:为|成)?\s*"
    r"(JSONL|JSON|CSV|XLSX|DOCX|PDF|HTML|MARKDOWN|MD|TXT|PPTX)",
    re.IGNORECASE,
)
_HARD_SOURCE_REQUIREMENT_PATTERN = re.compile(
    r"全部|所有|完整覆盖|不得遗漏|必须\s*(?:至少\s*)?\d+"
)
_COUNT_TOKEN = r"(?:\d{1,6}|[零〇一二两三四五六七八九十百千万]+)"
_RESULT_UNIT_TOKEN = (
    r"(?:家公司|家|项结果|条结果|个结果|个对象|个公司|项|条|份|篇|款|种|名)"
)
_RESULT_COUNT_WITH_UNIT_PATTERN = re.compile(
    rf"(?<!\d)({_COUNT_TOKEN})(?!\d)\s*{_RESULT_UNIT_TOKEN}"
)
_RESULT_COUNT_ONLY_PATTERN = re.compile(
    r"^(?:至少|最多|正好|约|大约|不少于|不超过|目标(?:为)?|需要)?\s*"
    rf"(?<!\d)({_COUNT_TOKEN})(?!\d)\s*$"
)
_HARD_QUANTITY_PATTERN = re.compile(r"至少|必须|不少于|不得少于|正好|恰好")
_UNSUPPORTED_COUNT_OPERATOR_PATTERN = re.compile(
    rf"(?:最多|至多|不超过|正好|恰好|大约|约)\s*"
    rf"(?<!\d){_COUNT_TOKEN}(?!\d)\s*{_RESULT_UNIT_TOKEN}|"
    rf"^(?:最多|至多|不超过|正好|恰好|大约|约)\s*"
    rf"(?<!\d){_COUNT_TOKEN}(?!\d)\s*$|"
    rf"(?<!\d){_COUNT_TOKEN}(?!\d)\s*{_RESULT_UNIT_TOKEN}\s*"
    r"(?:以内|以下|左右)|"
    rf"(?<!\d){_COUNT_TOKEN}(?!\d)\s*{_RESULT_UNIT_TOKEN}\s*"
    r"(?:上下|前后|整)"
    r"(?=$|[\s，。；、,;：:]|即可|就好|就行|完成|为准|足够|吧|为宜|"
    r"的(?:结果|范围|数量))"
)
_NEGATED_COMPLETENESS_PATTERN = re.compile(
    r"不(?:要求|需要)\s*(?:全部|所有|完整|全量)|无需\s*(?:全部|所有|完整|全量)|"
    r"允许遗漏|可以遗漏"
)
_TERMINAL = {
    "completed",
    "candidate_ready",
    "failed",
    "cancelled",
    "needs_input",
}


def _capability_catalog() -> CapabilityCatalog:
    """灰度入口显式执行纯新增目录迁移，不在普通任务启动时隐式写库。"""

    return CapabilityCatalog(
        SqliteCapabilityCatalogRepository(settings.webui_db_path)
    )


def _require_capability_gray(user: dict[str, Any]) -> None:
    if not is_admin_role(user.get("role")):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="任务级能力灰度只对管理员和超级管理员开放",
        )
    if not settings.pi_capability_host_enabled:
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail="任务级能力 Sidecar 灰度尚未启用",
        )


def _runtime_gate():
    """#13 冻结拦截共用的运行时门；与装载 Seam 是同一 check_mount 实现。"""
    from src.api.capability_governance_runtime import get_runtime_gate

    return get_runtime_gate()


def _runtime_gate_projection_governance():
    """选择列表的只读投影装配；不装配平台发布/签名依赖。"""
    from src.capability_governance import (
        CapabilityGovernance,
        SqliteCapabilityGovernanceRepository,
    )

    return CapabilityGovernance(
        _capability_catalog(),
        SqliteCapabilityGovernanceRepository(settings.webui_db_path),
    )


def _selectable_for_task(projection) -> bool:
    """新任务选择的三轴过滤；deprecated/revoked/quarantined/draft 不可选。

    历史冻结任务的恢复走 resolve_selection 路径，不受此谓词影响。
    """
    from src.capability_governance import (
        CapabilityEligibility,
        CapabilityLifecycle,
        CapabilityMaturity,
    )

    return (
        projection.maturity is CapabilityMaturity.VERIFIED
        and projection.lifecycle is CapabilityLifecycle.ACTIVE
        and projection.eligibility is CapabilityEligibility.ELIGIBLE
    )


def _check_freeze_gate(
    actor,
    pack_refs: tuple[CapabilityPackRef, ...],
    catalog: CapabilityCatalog,
    *,
    validation_target: CapabilityPackRef | None = None,
) -> None:
    """冻结前执行完整装载门 + 新任务可选谓词；拒绝以 409 失败关闭。

    装载门放行 DEPRECATED 是为历史恢复装载；冻结是「新任务」入口，
    必须再按三轴可选谓词拦截（AC3：deprecated 不进入新任务选择）。
    validation_target（#15 D9）：验证任务标记匹配的 ref 走豁免路径——
    仅跳过成熟度与可选谓词，门内其余条件（Owner/生命周期/资格）仍强制。
    """
    gate = _runtime_gate()
    governance = _runtime_gate_projection_governance()
    for ref in pack_refs:
        pack = catalog.resolve_pack(actor, ref.pack_id, ref.version)
        if pack is None:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="能力包不存在或当前用户不可见",
            )
        if pack.digest != ref.digest:
            # 身份失配是调用方输入错误（引用伪造或版本漂移），
            # 保持 422 语义；目录 freeze_selection 仍二次复核。
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail="能力包 digest 与引用不一致",
            )
        exempt = validation_target is not None and validation_target == ref
        try:
            gate.check_mount(actor, pack, validation_exempt=exempt)
        except CapabilityMountGateRejected as error:
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail=str(error),
            ) from error
        if not exempt and not _selectable_for_task(
            governance.runtime_projection_for_pack(pack)
        ):
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="能力当前不可用于新任务选择",
            )


def _inherit_capability_selection(
    owner_id: str,
    *,
    source_task_id: str,
    source_revision: int,
    target_task_id: str,
    target_revision: int,
) -> bool:
    """版本变化只复制冻结身份；没有目录表或旧选择时保持既有路径。"""

    resolver = DefaultCapabilityMounts(
        db_path=settings.webui_db_path,
        oci_layout_path=settings.capability_oci_layout_path,
        mount_root=settings.capability_mount_cache_path,
    )
    return resolver.copy_selection_for_owner(
        owner_id,
        source_task_id=source_task_id,
        source_revision=source_revision,
        target_task_id=target_task_id,
        target_revision=target_revision,
    )


class WorkspaceTaskCreateIn(BaseModel):
    model_config = ConfigDict(extra="forbid")

    objective_text: str = Field(min_length=1, max_length=20_000)
    upload_ids: tuple[str, ...] = ()
    source_snapshot_id: str | None = Field(default=None, min_length=1, max_length=160)
    must_include: tuple[str, ...] = ()
    explicit_exclusions: tuple[str, ...] = ()
    quantity_requirement: str | None = Field(default=None, min_length=1, max_length=500)
    completeness_requirement: str | None = Field(default=None, min_length=1, max_length=500)
    output_formats: tuple[str, ...] = ("xlsx",)
    table_output_contracts: tuple[TableOutputContract, ...] = ()
    provider: str = Field(default="local", min_length=1)
    model: str | None = Field(default=None, min_length=1)
    external_api_confirmed: bool = False
    runtime_version: RuntimeVersion = Field(
        default_factory=lambda: RuntimeVersion.LEGACY,
    )
    permission_profile: PermissionProfile = PermissionProfile.STANDARD
    model_connection_id: str | None = Field(
        default=None,
        min_length=1,
        max_length=100,
    )
    model_connection_model: str | None = Field(default=None, min_length=1, max_length=200)
    capability_pack_refs: tuple[CapabilityPackRef, ...] = ()
    # #15 D9 验证任务标记：本任务是为验证该个人 draft 能力而创建；
    # 仅在 create_task 校验后随冻结 selection 落库。
    validation_target: CapabilityPackRef | None = None
    context_purpose: str = Field(default="web_research", min_length=1, max_length=80)
    context_selection: TaskContextSelection | None = None
    context_preview_sha256: str | None = Field(
        default=None,
        pattern=r"^sha256:[0-9a-f]{64}$",
    )

    @field_validator("objective_text", "provider")
    @classmethod
    def strip_text(cls, value: str) -> str:
        return value.strip()

    @field_validator("upload_ids")
    @classmethod
    def unique_uploads(cls, value: tuple[str, ...]) -> tuple[str, ...]:
        if len(value) != len(set(value)):
            raise ValueError("upload_ids 不得重复")
        return value

    @field_validator("must_include", "explicit_exclusions")
    @classmethod
    def bounded_web_constraints(
        cls, value: tuple[str, ...]
    ) -> tuple[str, ...]:
        normalized = tuple(item.strip() for item in value)
        if len(normalized) > 50:
            raise ValueError("网页目标合同每类最多填写 50 条约束")
        if any(not item or len(item) > 500 for item in normalized):
            raise ValueError("网页目标合同每条约束长度必须为 1 至 500 个字符")
        if len(normalized) != len(set(normalized)):
            raise ValueError("网页目标合同约束不得重复")
        return normalized

    @field_validator("output_formats")
    @classmethod
    def valid_formats(cls, value: tuple[str, ...]) -> tuple[str, ...]:
        normalized = tuple(item.strip().lower() for item in value)
        if not normalized:
            raise ValueError("至少选择一种正式输出格式")
        invalid = set(normalized) - _FORMATS
        if invalid:
            raise ValueError(f"不支持的正式输出格式：{sorted(invalid)}")
        if len(normalized) != len(set(normalized)):
            raise ValueError("output_formats 不得重复")
        return normalized

    @field_validator("capability_pack_refs")
    @classmethod
    def unique_capabilities(
        cls,
        value: tuple[CapabilityPackRef, ...],
    ) -> tuple[CapabilityPackRef, ...]:
        identities = {(item.pack_id, item.version) for item in value}
        if len(identities) != len(value):
            raise ValueError("capability_pack_refs 不得重复")
        return value

    @model_validator(mode="after")
    def validate_validation_target(self) -> "WorkspaceTaskCreateIn":
        if bool(self.context_selection) != bool(self.context_preview_sha256):
            raise ValueError("上下文选择和已确认预览哈希必须同时提供")
        if self.context_selection is not None and self.source_snapshot_id is None:
            raise ValueError("P1-01 上下文复用当前只开放给网页任务")
        if bool(self.upload_ids) == bool(self.source_snapshot_id):
            raise ValueError("任务必须且只能选择上传文件或一个网页来源快照")
        web_contract_values = (
            self.must_include,
            self.explicit_exclusions,
            self.quantity_requirement,
            self.completeness_requirement,
        )
        if self.source_snapshot_id is not None:
            if self.quantity_requirement is None or self.completeness_requirement is None:
                raise ValueError("网页任务启动前必须确认数量和完整性要求")
            if self.runtime_version is not RuntimeVersion.PI:
                raise ValueError("网页来源当前只能由 AgentKernel 统一运行时执行")
            constraint_lines = (
                *(f"必须包含：{item}" for item in self.must_include),
                *(f"明确不要：{item}" for item in self.explicit_exclusions),
                f"数量要求：{self.quantity_requirement}",
                f"完整性边界：{self.completeness_requirement}",
            )
            execution_objective = "\n".join(
                (self.objective_text, "", "执行约束：", *constraint_lines)
            )
            if len(execution_objective) > 20_000:
                raise ValueError("网页目标与执行约束合计不能超过 20000 个字符")
        elif any(web_contract_values):
            raise ValueError("网页目标合同只能绑定网页来源快照")
        contract_formats = [
            item.format for item in self.table_output_contracts
        ]
        if len(contract_formats) != len(set(contract_formats)):
            raise ValueError("同一输出格式只能冻结一个表格契约")
        if not set(contract_formats).issubset(self.output_formats):
            raise ValueError("表格输出契约必须绑定正式输出格式")
        if (
            self.table_output_contracts
            and "runtime_version" in self.model_fields_set
            and self.runtime_version is not RuntimeVersion.PI
        ):
            raise ValueError("表格输出契约只能由 Pi Runtime 执行")
        # #15 D9：验证目标必须同时出现在能力选择中（否则豁免无载体）。
        if (
            self.validation_target is not None
            and not any(
                item == self.validation_target for item in self.capability_pack_refs
            )
        ):
            raise ValueError("验证目标必须同时出现在能力选择中")
        return self


class WorkspaceAnswerIn(BaseModel):
    model_config = ConfigDict(extra="forbid")

    answer: str = Field(min_length=1, max_length=10_000)


class WorkspaceTurnIn(BaseModel):
    """运行中追问或修改；先理解差异，不直接改写任务。"""

    model_config = ConfigDict(extra="forbid")

    text: str = Field(min_length=1, max_length=20_000)
    external_api_confirmed: bool = False

    @field_validator("text")
    @classmethod
    def strip_text(cls, value: str) -> str:
        return value.strip()


class CandidateReverificationIn(BaseModel):
    model_config = ConfigDict(extra="forbid")

    expected_revision: int = Field(ge=1)
    expected_previous_attempt_id: str = Field(min_length=1, max_length=160)
    expected_candidate_set_hash: str | None = Field(
        default=None,
        pattern=r"^[0-9a-f]{64}$",
    )
    expected_target_ruleset_hash: str | None = Field(
        default=None,
        pattern=r"^[0-9a-f]{64}$",
    )
    legacy_ruleset_unknown_acknowledged: Literal[True] | None = None
    authorization_text_version: Literal["legacy-rebaseline-v1"] | None = None
    external_api_confirmed: bool = False
    accept_duplicate_provider_cost: bool = False
    historical_authority_recovery: (
        HistoricalAuthorityRecoveryConfirmation | None
    ) = None


class CandidateVerificationPublishIn(BaseModel):
    model_config = ConfigDict(extra="forbid")

    expected_revision: int = Field(ge=1)


class WorkspaceRevisionDecisionIn(BaseModel):
    model_config = ConfigDict(extra="forbid")

    mode: RevisionSwitchMode
    external_api_confirmed: bool = False


class WorkspaceRevisionIn(BaseModel):
    model_config = ConfigDict(extra="forbid")

    instruction: str = Field(min_length=1, max_length=20_000)
    output_formats: tuple[str, ...] | None = None
    table_output_contracts: tuple[TableOutputContract, ...] | None = None
    external_api_confirmed: bool = False
    expected_active_revision: int = Field(ge=1)
    source_snapshot_id: str | None = Field(default=None, min_length=1, max_length=160)
    accepted_candidate_set_hash: str | None = Field(
        default=None,
        pattern=r"^[0-9a-f]{64}$",
    )
    accepted_result_count: int | None = Field(default=None, ge=1)

    @field_validator("instruction")
    @classmethod
    def strip_instruction(cls, value: str) -> str:
        return value.strip()

    @field_validator("output_formats")
    @classmethod
    def valid_revision_formats(
        cls, value: tuple[str, ...] | None
    ) -> tuple[str, ...] | None:
        if value is None:
            return None
        normalized = tuple(item.strip().lower() for item in value)
        if not normalized:
            raise ValueError("至少选择一种正式输出格式")
        invalid = set(normalized) - _FORMATS
        if invalid:
            raise ValueError(f"不支持的正式输出格式：{sorted(invalid)}")
        return normalized

    @model_validator(mode="after")
    def validate_table_output_contracts(self) -> "WorkspaceRevisionIn":
        if bool(self.accepted_candidate_set_hash) != bool(self.accepted_result_count):
            raise ValueError("接受缺口必须同时绑定 Candidate 身份和实际结果数")
        if self.table_output_contracts is None:
            return self
        contract_formats = [
            item.format for item in self.table_output_contracts
        ]
        if len(contract_formats) != len(set(contract_formats)):
            raise ValueError("同一输出格式只能冻结一个表格契约")
        if (
            self.output_formats is not None
            and not set(contract_formats).issubset(self.output_formats)
        ):
            raise ValueError("表格输出契约必须绑定正式输出格式")
        return self


class WorkspaceContextPreviewIn(BaseModel):
    model_config = ConfigDict(extra="forbid")

    purpose: str = Field(default="web_research", min_length=1, max_length=80)
    objective_text: str = Field(min_length=1, max_length=20_000)
    output_formats: tuple[str, ...] = Field(min_length=1)
    selection: TaskContextSelection = Field(default_factory=TaskContextSelection)

    @field_validator("objective_text")
    @classmethod
    def strip_objective(cls, value: str) -> str:
        return value.strip()

    @field_validator("output_formats")
    @classmethod
    def valid_context_formats(cls, value: tuple[str, ...]) -> tuple[str, ...]:
        normalized = tuple(item.strip().lower() for item in value)
        invalid = set(normalized) - _FORMATS
        if invalid:
            raise ValueError(f"不支持的正式输出格式：{sorted(invalid)}")
        return normalized


class CandidateGapActionIn(BaseModel):
    model_config = ConfigDict(extra="forbid")

    action: Literal[
        "accept_gap",
        "reject_gap",
        "supplement_source",
        "refresh_source",
    ]
    expected_revision: int = Field(ge=1)
    expected_candidate_set_hash: str = Field(pattern=r"^[0-9a-f]{64}$")
    external_api_confirmed: bool = False


class SourceRefreshIn(BaseModel):
    model_config = ConfigDict(extra="forbid")

    expected_active_revision: int = Field(ge=1)
    external_api_confirmed: bool = False
    resume_unknown: bool = False


def _uploads() -> UploadStore:
    return UploadStore(
        root=settings.data_prep_upload_root,
        max_bytes=settings.data_prep_max_upload_bytes,
    )


def _runtime_repository() -> AgenticRuntimeRepository:
    return AgenticRuntimeRepository(settings.webui_db_path)


def _source_acquisition_service() -> SourceAcquisitionService:
    return SourceAcquisitionService(
        SourceAcquisitionRepository(settings.webui_db_path),
        AnonymousWebFetcher(),
    )


def _task_context_service() -> TaskContextService:
    return TaskContextService(TaskContextRepository(settings.webui_db_path))


def _inherit_task_context_hook(
    base_hook: Callable[[sqlite3.Connection], None] | None,
    *,
    owner_id: str,
    source_task_id: str,
    source_revision: int,
    target_task_id: str,
    target_revision: int,
    objective_text: str,
    output_formats: tuple[str, ...],
) -> Callable[[sqlite3.Connection], None] | None:
    """把既有精确引用带到已确认的新 Revision，不重新读取浮动目录。"""

    service = _task_context_service()
    preview = service.carry_forward(
        owner_id=owner_id,
        source_task_id=source_task_id,
        source_revision=source_revision,
        target_task_id=target_task_id,
        target_revision=target_revision,
        objective_text=objective_text,
        output_formats=output_formats,
    )
    if preview is None:
        return base_hook

    def bind_context(connection: sqlite3.Connection) -> None:
        if base_hook is not None:
            base_hook(connection)
        service.freeze(
            connection,
            owner_id=owner_id,
            task_id=target_task_id,
            revision=target_revision,
            preview=preview,
            expected_preview_sha256=preview.preview_sha256,
        )

    return bind_context


@router.get("/context-options")
def list_context_options(
    purpose: str = Query(default="web_research", min_length=1, max_length=80),
    user=Depends(get_current_user),
):
    """只列出当前 Owner、当前用途可显式选择的模板与记忆摘要。"""

    options = _task_context_service().options(user["user_id"], purpose)
    return {
        key: [item.model_dump(mode="json") for item in items]
        for key, items in options.items()
    }


@router.post("/context-preview")
def preview_task_context(
    payload: WorkspaceContextPreviewIn,
    user=Depends(get_current_user),
):
    """启动前生成可检查草案；本接口不创建或修改 TaskRevision。"""

    try:
        preview = _task_context_service().preview(
            owner_id=user["user_id"],
            purpose=payload.purpose,
            objective_text=payload.objective_text,
            output_formats=payload.output_formats,
            selection=payload.selection,
        )
    except KeyError as exc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)) from exc
    except ValueError as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail=str(exc),
        ) from exc
    return preview.model_dump(mode="json")


def _rollout_actor(user: dict[str, Any]) -> RolloutActor:
    role = (
        "super_admin"
        if user.get("role") == "super_admin"
        else "admin"
        if is_admin_role(user.get("role"))
        else "user"
    )
    return RolloutActor(actor_id=user["user_id"], role=role)


@dataclass(frozen=True)
class _RevisionRoutingPlan:
    selected_runtime: RuntimeVersion
    repository: SqliteRuntimeRoutingRepository | None = None
    task_revision: RuntimeTaskRevisionRef | None = None
    actor: RolloutActor | None = None
    rollout: RolloutSnapshot | None = None


def _preview_revision_runtime(
    user: dict[str, Any],
    *,
    task_id: str,
    revision: int,
    requested_runtime: RuntimeVersion,
) -> _RevisionRoutingPlan:
    repository = open_runtime_routing_repository(settings.webui_db_path)
    if repository is None:
        return _RevisionRoutingPlan(selected_runtime=requested_runtime)
    routing = RuntimeRouting(repository)
    task_revision = RuntimeTaskRevisionRef(
        owner_id=user["user_id"],
        task_id=task_id,
        revision=revision,
        requested_runtime=requested_runtime,
    )
    actor = _rollout_actor(user)
    selected_runtime, rollout = routing.preview(task_revision, actor)
    return _RevisionRoutingPlan(
        selected_runtime=selected_runtime,
        repository=repository,
        task_revision=task_revision,
        actor=actor,
        rollout=rollout,
    )


def _prepare_runtime_binding(
    plan: _RevisionRoutingPlan,
    runtime_config: RuntimeTaskConfig,
) -> tuple[
    RuntimeVersion,
    Callable[[sqlite3.Connection], None] | None,
]:
    assignment = None
    if plan.repository is not None:
        assert plan.task_revision is not None
        assert plan.actor is not None
        assert plan.rollout is not None
        assignment = RuntimeAssignment(
            task_revision=plan.task_revision,
            runtime_version=plan.selected_runtime,
            rollout_mode=plan.rollout.mode,
            gate_snapshot_id=plan.rollout.active_gate_snapshot_id,
            assigned_by=plan.actor.actor_id,
            assigned_at=datetime.now(timezone.utc),
        )
    runtime_repository = _runtime_repository()

    def bind_runtime(connection: sqlite3.Connection) -> None:
        try:
            if assignment is not None:
                assert plan.repository is not None
                plan.repository.create_assignment_in_transaction(
                    connection,
                    assignment,
                )
            runtime_repository.register_in_transaction(connection, runtime_config)
        except ValueError as exc:
            raise RuntimeError(str(exc)) from exc

    return plan.selected_runtime, bind_runtime


def _steering_repository() -> SqliteSteeringRepository:
    return SqliteSteeringRepository(settings.webui_db_path)


def _public_runtime(
    user_id: str,
    task_id: str,
    revision: int,
) -> dict[str, Any]:
    repository = _runtime_repository()
    row = repository.get(user_id, task_id, revision)
    if row is None:
        return {
            "runtime_version": RuntimeVersion.LEGACY.value,
            "permission_profile": PermissionProfile.STANDARD.value,
            "model_connection_id": None,
            "external_api_confirmed": False,
            "status": None,
            "candidates": [],
            "provider_usage": [],
            "events": [],
        }
    coverage = (
        repository.get_coverage(
            user_id=user_id,
            task_id=task_id,
            revision=revision,
            run_id=row["run_id"],
        )
        if row["run_id"]
        else None
    )
    candidate_visible = row["status"].value == "candidate_ready"
    verification = (
        row["verification"].model_dump(mode="json")
        if candidate_visible and row["verification"]
        else None
    )
    candidate_coverage = (
        repository.get_candidate_coverage(
            user_id=user_id,
            task_id=task_id,
            revision=revision,
            candidate_set_hash=str(row["verified_candidate_set_hash"]),
        )
        if row.get("verified_candidate_set_hash")
        else None
    )
    if verification and verification.get("status") == "inconclusive":
        for check in verification.get("checks", []):
            if check.get("code") == "semantic_goal" and not check.get("passed"):
                # 兼容历史任务：持久化的第三方解析栈不得继续暴露给普通工作台。
                check["summary"] = (
                    "语义验证服务暂时不可用，请稍后重新验证候选。"
                )
    payload = {
        "runtime_version": row["runtime_version"].value,
        "permission_profile": row["permission_profile"].value,
            "model_connection_id": row["model_connection_id"],
            "model_connection_model": row["model_connection_model"],
        "external_api_confirmed": row["external_api_confirmed"],
        "status": row["status"].value,
        "run_id": row["run_id"],
        "session_file": row["session_file"],
        "summary": (
            (row["request"] or {}).get("objective_text")
            if row["request"]
            else None
        ),
        "candidates": [
            {
                **item.public_dict(task_id=task_id, revision=revision),
                "download_allowed": _candidate_download_allowed(
                    row["verification"],
                    candidate_format=item.format,
                ),
            }
            for item in (row["candidates"] if candidate_visible else ())
        ],
        "verification": verification,
        "candidate_coverage": (
            candidate_coverage.model_dump(mode="json")
            if candidate_coverage is not None
            else None
        ),
        "gap_actions": repository.list_gap_actions(
            user_id=user_id,
            task_id=task_id,
            source_revision=revision,
        ),
        "failure": row["failure"],
        "provider_usage": (
            get_default_broker().list_usage(
                user_id,
                task_id=task_id,
                revision=revision,
            )
            if row["model_connection_id"]
            else []
        ),
        "events": repository.list_events(
            user_id, task_id, revision
        ),
        "coverage": (
            {
                "contract": coverage[0].model_dump(mode="json"),
                "ledger": coverage[1].model_dump(mode="json"),
                "progress": coverage[1].public_progress(),
            }
            if coverage
            else None
        ),
    }
    if candidate_visible:
        try:
            offer = (
                get_semantic_workspace_manager().inspect_candidate_reverification(
                    user_id,
                    task_id,
                    revision,
                )
            )
        except ReverificationContractError:
            history = SqliteCandidateVerificationRepository(
                settings.webui_db_path
            ).list_for_candidate(
                user_id,
                task_id=task_id,
                revision=revision,
                run_id=str(row["run_id"]),
                candidate_set_hash=str(row["verified_candidate_set_hash"]),
            )
            latest_attempt = history[-1] if history else None
            if (
                latest_attempt is None
                or latest_attempt.ruleset_identity_status
                is not RulesetIdentityStatus.LEGACY_UNVERSIONED
            ):
                # 只有可证明的旧版 Attempt 可以降级读取；现代记录损坏必须显式失败。
                raise
            # 旧任务可以继续读取，但缺失的冻结上下文不能被推断或补写。
            payload["latest_verification_attempt"] = {
                "attempt_id": latest_attempt.attempt_id,
                "status": latest_attempt.status.value,
                "reason": latest_attempt.reason_code.value,
                "ruleset_identity_status": (
                    latest_attempt.ruleset_identity_status.value
                ),
            }
            payload["reverification_offer"] = None
            payload["reverification_unavailable_reason"] = (
                "该历史任务缺少可证明的冻结运行信息，暂不能重新验证。"
            )
            return payload
        payload["latest_verification_attempt"] = {
            "attempt_id": offer.previous_attempt_id,
            "status": (
                offer.previous_status.value if offer.previous_status else None
            ),
            "reason": (
                offer.previous_reason.value if offer.previous_reason else None
            ),
            "ruleset_identity_status": (
                offer.ruleset_identity_status.value
                if offer.ruleset_identity_status
                else None
            ),
        }
        payload["reverification_offer"] = {
            "eligible": offer.eligible,
            "reason": offer.reason.value if offer.reason else None,
            "blockers": list(offer.blockers),
            "ruleset_changed": offer.ruleset_changed,
            "ruleset_change_summary": offer.ruleset_change_summary,
            "requires_provider": offer.requires_provider,
            "connection_id": offer.connection_id,
            "model_id": offer.model_id,
            "candidate_count": offer.candidate_count,
            "candidate_formats": list(offer.candidate_formats),
            "candidate_set_hash": offer.candidate_set_hash,
            "target_ruleset_hash": offer.target_ruleset_hash,
            "egress_categories": list(offer.egress_categories),
            "egress_summary": offer.egress_summary,
            "historical_authority_recovery": (
                offer.historical_authority_recovery.model_dump(mode="json")
                if offer.historical_authority_recovery is not None
                else None
            ),
        }
        payload["awaiting_publication"] = bool(
            offer.awaiting_publication
            and not (
                candidate_coverage is not None
                and candidate_coverage.is_partial
            )
        )
    return payload


_SAFE_TEXT_CANDIDATE_FORMATS = {"csv", "json", "jsonl", "markdown", "txt"}


def _candidate_download_allowed(
    verification,
    *,
    candidate_format: str,
) -> bool:
    """安全文本候选允许 Owner 人工检查；主动内容和复杂容器继续失败关闭。"""

    if candidate_format in _SAFE_TEXT_CANDIDATE_FORMATS:
        # 候选进入仓库前已经过格式重开、路径归属和哈希检查。来源证据失败只表示
        # 内容不能成为正式交付，不应让任务所有者失去检查自己纯文本候选的能力。
        return True

    if verification is None:
        return True
    blocked_codes = {"manifest", "artifact_set", "source_grounding"}
    return not any(
        not check.passed and check.code in blocked_codes
        for check in verification.checks
    )


def _freeze_goal_contract(
    payload: WorkspaceTaskCreateIn,
    *,
    objective: str,
    source_snapshot: dict[str, Any],
) -> dict[str, Any]:
    """把用户确认的结果门槛和获准来源范围写成显式 GoalContract。"""

    quantity = payload.quantity_requirement or ""
    completeness = payload.completeness_requirement or ""
    confirmed_requirements = f"{quantity}\n{completeness}"
    if _UNSUPPORTED_COUNT_OPERATOR_PATTERN.search(quantity):
        # 当前覆盖契约只表达最低目标，不能丢掉“最多/正好/约”等运算符后
        # 全部按“至少 N 项”执行。在运算符进入领域模型前必须失败关闭。
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail=(
                "当前暂不支持结果数量上限、精确值或约数，请改为"
                "“至少 10 家”或“尽可能多”"
            ),
        )
    # 数量只从用户单独确认的数量字段提取，并要求结果单位或纯数量表达；
    # “最近 30 天”等时间窗口不能悄悄变成 30 项硬门槛。
    count_match = _RESULT_COUNT_WITH_UNIT_PATTERN.search(quantity)
    if count_match is None:
        count_match = _RESULT_COUNT_ONLY_PATTERN.fullmatch(quantity.strip())
    require_all = bool(
        re.search(r"全部|所有|不得遗漏|完整|必须全量", confirmed_requirements)
    ) and not bool(_NEGATED_COMPLETENESS_PATTERN.search(confirmed_requirements))
    target_count = _parse_result_count(count_match.group(1)) if count_match else None
    if target_count is None and _HARD_QUANTITY_PATTERN.search(quantity):
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail="无法确定硬性结果数量，请使用“至少 10 家公司”等明确数量重新确认",
        )
    strict = bool(require_all or target_count is not None)
    return {
        "objective": objective,
        "must_include": list(payload.must_include),
        "explicit_exclusions": list(payload.explicit_exclusions),
        "quantity_requirement": payload.quantity_requirement,
        "completeness_requirement": payload.completeness_requirement,
        "coverage": {
            "strictness": "strict" if strict else "exploratory",
            "target_result_count": target_count,
            "require_all": require_all,
            "exploratory_target": None if strict else confirmed_requirements.strip(),
            "authorized_source_scope": source_snapshot["allowed_scope"],
        },
    }


def _parse_result_count(value: str) -> int | None:
    """解析已绑定结果单位的阿拉伯或常用中文整数；越界时失败关闭。"""

    if value.isdigit():
        parsed = int(value)
        return parsed if 1 <= parsed <= 1_000_000 else None
    digits = {"零": 0, "〇": 0, "一": 1, "二": 2, "两": 2, "三": 3,
              "四": 4, "五": 5, "六": 6, "七": 7, "八": 8, "九": 9}
    units = {"十": 10, "百": 100, "千": 1000}
    total = 0
    section = 0
    number = 0
    for char in value:
        if char in digits:
            number = digits[char]
        elif char in units:
            section += (number or 1) * units[char]
            number = 0
        elif char == "万":
            total += (section + number or 1) * 10_000
            section = 0
            number = 0
        else:
            return None
    parsed = total + section + number
    return parsed if 1 <= parsed <= 1_000_000 else None


@router.get("/guidance")
def workspace_guidance(user=Depends(get_current_user)):
    del user
    path = (
        Path(__file__).resolve().parents[3]
        / "config"
        / "workspace_guidance.json"
    )
    return json.loads(path.read_text(encoding="utf-8"))


def _task_or_404(user_id: str, task_id: str) -> dict[str, Any]:
    task = get_store().get_semantic_workspace_task(user_id, task_id)
    if task is None:
        raise HTTPException(status_code=404, detail="工作台任务不存在")
    return task


@router.get("/capabilities")
def list_gray_capabilities(user=Depends(get_current_user)):
    """只列出已验证的平台能力；草稿和执行配置不进入灰度选择界面。"""

    _require_capability_gray(user)
    actor = catalog_actor_from_user(user)
    catalog = _capability_catalog()
    governance = _runtime_gate_projection_governance()
    items = []
    for pack in catalog.list_visible_packs(actor):
        if pack.maturity is not CapabilityMaturity.VERIFIED:
            continue
        projection = governance.runtime_projection_for_pack(pack)
        if not _selectable_for_task(projection):
            # deprecated/revoked/quarantined 不进入新任务选择（AC3）。
            continue
        manifest = dict(pack.manifest)
        kind = manifest.get("kind", "capability_pack")
        if kind not in {
            "tool",
            "mcp_local",
            "skill",
            "dependency_bundle",
            "capability_pack",
        }:
            kind = "capability_pack"
        items.append({
            "pack_id": pack.pack_id,
            "version": pack.version,
            "digest": pack.digest,
            "name": manifest.get("display_name") or pack.pack_id,
            "kind": kind,
            "purpose": (
                manifest.get("purpose")
                or "提供当前任务所需的专业处理能力"
            ),
            "scope": pack.scope.value,
            # 推荐指针（#14 回滚命令折叠）；推荐是默认值不是约束。
            "recommended": (
                projection.recommended_version == pack.version
            ),
        })
    # 推荐版本置顶；无指针的 pack 保持目录既有顺序（旧路径零回归）。
    items.sort(key=lambda item: not item["recommended"])
    return {"enabled": True, "items": items}


@router.get("/tasks/{task_id}/capabilities")
def get_task_capabilities(
    task_id: str,
    revision: int = Query(default=1, ge=1),
    user=Depends(get_current_user),
):
    _task_or_404(user["user_id"], task_id)
    resolver = DefaultCapabilityMounts(
        db_path=settings.webui_db_path,
        oci_layout_path=settings.capability_oci_layout_path,
        mount_root=settings.capability_mount_cache_path,
    )
    items = resolver.describe_for_owner(user["user_id"], task_id, revision)
    return {
        "task_id": task_id,
        "revision": revision,
        "items": [item.model_dump(mode="json") for item in items],
    }


def _revision_events(
    events: list[dict[str, Any]],
    revision: int,
) -> list[dict[str, Any]]:
    """按 revision_created 边界切分工作台事件，避免历史版本串入后续轨迹。"""
    boundaries = [
        event["sequence"]
        for event in events
        if event["event_type"] == "revision_created"
    ]
    if revision >= 2 and len(boundaries) < revision - 1:
        return []
    start = boundaries[revision - 2] if revision >= 2 else 0
    end = (
        boundaries[revision - 1]
        if revision - 1 < len(boundaries)
        else None
    )
    return [
        event
        for event in events
        if event["sequence"] >= start
        and (end is None or event["sequence"] < end)
    ]


_PROGRESS_STAGE_ALIASES = {
    "queued": ProgressStage.UNDERSTAND,
    "interpret": ProgressStage.UNDERSTAND,
    "understand": ProgressStage.UNDERSTAND,
    "goal_interpretation": ProgressStage.UNDERSTAND,
    "inspect": ProgressStage.INSPECT_SOURCES,
    "inspect_sources": ProgressStage.INSPECT_SOURCES,
    "source_probe": ProgressStage.INSPECT_SOURCES,
    "source_discovery": ProgressStage.INSPECT_SOURCES,
    "bind": ProgressStage.INSPECT_SOURCES,
    "prepare_capabilities": ProgressStage.PREPARE_CAPABILITIES,
    "execute": ProgressStage.EXECUTE,
    "repair": ProgressStage.EXECUTE,
    "agent_execute": ProgressStage.EXECUTE,
    "evidence_read": ProgressStage.EXECUTE,
    "verify": ProgressStage.VERIFY,
    "check": ProgressStage.VERIFY,
    "verify_coverage": ProgressStage.VERIFY,
    "deliver": ProgressStage.DELIVER,
    "output": ProgressStage.DELIVER,
}

_PROGRESS_TOOL_SUMMARIES = {
    "freeze_coverage": {
        "tool.started": "正在确认任务范围与完成条件",
        "tool.completed": "已确认任务范围与完成条件",
        "tool.failed": "任务范围确认遇到问题，正在调整",
    },
    "inspect_source": {
        "tool.started": "正在检查来源结构",
        "tool.completed": "已完成来源结构检查",
        "tool.failed": "来源结构检查遇到问题，正在调整",
    },
    "discover_content": {
        "tool.started": "正在定位候选内容",
        "tool.completed": "已完成候选内容定位",
        "tool.failed": "候选内容定位遇到问题，正在调整",
    },
    "read_evidence": {
        "tool.started": "正在读取并核对来源证据",
        "tool.completed": "已完成来源证据读取",
        "tool.failed": "来源证据读取遇到问题，正在调整",
    },
    "propose_completion": {
        "tool.started": "正在检查结果完整性",
        "tool.completed": "已完成结果完整性检查",
        "tool.failed": "结果完整性检查未通过，正在补充处理",
    },
}


def _progress_stage(value: str) -> ProgressStage:
    normalized = value.strip().lower()
    return _PROGRESS_STAGE_ALIASES.get(normalized, ProgressStage.EXECUTE)


def _progress_summary(event: dict[str, Any], details: dict[str, Any]) -> str:
    tool = str(details.get("tool") or "")
    runtime_event_type = str(
        details.get("runtime_event_type")
        or event.get("event_type")
        or event.get("type")
        or ""
    )
    return _PROGRESS_TOOL_SUMMARIES.get(tool, {}).get(
        runtime_event_type,
        str(event.get("summary") or "正在处理"),
    )


def _structured_progress_events(task: dict[str, Any]) -> tuple[StructuredProgressEvent, ...]:
    projected: list[StructuredProgressEvent] = []
    raw_events = list(task["events"])
    raw_events.extend(
        {
            **event,
            "stage": event.get("node", "execute"),
        }
        for event in task["harness_events"]
    )
    for index, event in enumerate(raw_events, start=1):
        details = event.get("details") or {}
        raw_progress = details.get("progress")
        progress = None
        if isinstance(raw_progress, dict):
            try:
                progress = ProgressValue.model_validate(raw_progress)
            except ValueError:
                progress = None
        try:
            audience = ProgressAudience(details.get("audience", "all"))
        except ValueError:
            audience = ProgressAudience.ALL
        projected.append(
            StructuredProgressEvent(
                event_id=event.get("event_id") or f"legacy_{index}",
                # 工作台与 Harness 各自维护 sequence；兼容投影按合并后的
                # 事实顺序重新编号，避免两个 sequence=1 造成活动阶段倒退。
                sequence=index,
                task_id=task["task_id"],
                revision=int(task["viewing_revision"]),
                run_id=event.get("run_id") or task.get("run_id"),
                stage=_progress_stage(str(event.get("stage") or "execute")),
                event_type=str(event.get("event_type") or event.get("type") or "progress"),
                # 主时间线对新旧任务统一使用业务语言；原始事件仍供管理员诊断。
                summary=_progress_summary(event, details),
                progress=progress,
                refs=(details.get("refs") if isinstance(details.get("refs"), dict) else {}),
                action=(details.get("action") if isinstance(details.get("action"), dict) else None),
                audience=audience,
                created_at=event.get("created_at") or datetime.now().astimezone(),
            )
        )
    return tuple(projected)


def _task_detail(
    user_id: str,
    task_id: str,
    *,
    revision: int | None = None,
    audience: ProgressAudience = ProgressAudience.USER,
) -> dict[str, Any]:
    store = get_store()
    task = _task_or_404(user_id, task_id)
    task["current_status"] = task["status"]
    task["current_revision"] = task["active_revision"]
    task["revisions"] = store.list_semantic_workspace_revisions(
        user_id, task_id
    )
    selected_revision = (
        store.get_semantic_workspace_revision(
            user_id, task_id, revision
        )
        if revision is not None
        else store.get_semantic_workspace_revision(
            user_id, task_id, task["active_revision"]
        )
    )
    if selected_revision is None:
        raise HTTPException(status_code=404, detail="结果版本不存在")
    task["viewing_revision"] = selected_revision["revision"]
    if selected_revision["revision"] != task["active_revision"]:
        task.update(
            {
                key: selected_revision[key]
                for key in (
                    "objective_text",
                    "output_formats",
                    "plan_id",
                    "logical_revision",
                    "binding_revision",
                    "run_id",
                    "status",
                    "summary",
                )
            }
        )
    task["events"] = _revision_events(
        store.list_semantic_workspace_events(user_id, task_id),
        selected_revision["revision"],
    )
    task["uploads"] = []
    for upload_id in task["upload_ids"]:
        try:
            item = _uploads().resolve(user_id, upload_id)
        except PermissionError:
            continue
        task["uploads"].append(
            item.model_dump(exclude={"user_id", "storage_path"})
        )
    task["plan"] = None
    if task["plan_id"] and task["logical_revision"]:
        plan = store.get_semantic_plan_revision(
            user_id,
            task["plan_id"],
            task["logical_revision"],
        )
        if plan:
            task["plan"] = {
                "revision": plan["revision"],
                "summary": plan["summary"],
                "plan": plan["plan"],
                "diagnostics": plan["diagnostics"],
            }
    task["run"] = None
    task["attempts"] = []
    task["harness_events"] = []
    task["delivery"] = None
    if task["run_id"]:
        task["run"] = store.get_semantic_harness_run(
            user_id, task["run_id"]
        )
        task["attempts"] = store.list_semantic_harness_attempts(
            user_id, task["run_id"]
        )
        task["harness_events"] = store.list_semantic_harness_events(
            user_id, task["run_id"]
        )
        task["delivery"] = store.latest_semantic_delivery(
            user_id, task["run_id"]
        )
    task["agentic_runtime"] = _public_runtime(
        user_id,
        task_id,
        selected_revision["revision"],
    )
    task["runtime_version"] = task["agentic_runtime"]["runtime_version"]
    task["permission_profile"] = task["agentic_runtime"][
        "permission_profile"
    ]
    task["model_connection_id"] = task["agentic_runtime"][
        "model_connection_id"
    ]
    task["web_source"] = None
    frozen_context = TaskContextRepository(
        settings.webui_db_path
    ).get_frozen(user_id, task_id, int(selected_revision["revision"]))
    task["task_context"] = (
        frozen_context.model_dump(mode="json")
        if frozen_context is not None
        else None
    )
    web_contract = store.get_web_task_contract(
        user_id,
        task_id,
        selected_revision["revision"],
    )
    if web_contract is not None:
        snapshot = SourceAcquisitionRepository(
            settings.webui_db_path
        ).get_snapshot(user_id, web_contract["source_snapshot_id"])
        if snapshot is None:
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="冻结网页来源快照已缺失",
            )
        task["web_source"] = {**web_contract, "snapshot": snapshot}
    task["progress"] = ProgressProjection().project(
        _structured_progress_events(task),
        audience=audience,
        task_status=task["status"],
    ).model_dump(mode="json")
    return task


@router.post("/tasks", status_code=status.HTTP_202_ACCEPTED)
async def create_task(
    payload: WorkspaceTaskCreateIn,
    idempotency_key: str | None = Header(
        default=None,
        alias="Idempotency-Key",
    ),
    user=Depends(get_current_user),
):
    user_id = user["user_id"]
    user_objective = payload.objective_text
    idempotency_payload = payload.model_dump(
        mode="json",
        exclude_unset=True,
    )
    task_id = f"workspace_{uuid.uuid4().hex[:16]}"
    requested_runtime = (
        payload.runtime_version
        if "runtime_version" in payload.model_fields_set
        else None
    )
    routing_repository = open_runtime_routing_repository(settings.webui_db_path)
    routing_ref = None
    rollout_actor = None
    routing_preview = None
    if routing_repository is not None:
        routing_ref = RuntimeTaskRevisionRef(
            owner_id=user_id,
            task_id=task_id,
            revision=1,
            requested_runtime=requested_runtime,
        )
        rollout_actor = _rollout_actor(user)
        selected_runtime, routing_preview = RuntimeRouting(
            routing_repository
        ).preview(routing_ref, rollout_actor)
        if routing_preview.p0_blocked and requested_runtime is RuntimeVersion.PI:
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="vNext 因 P0 回归已暂停；请创建 Legacy 任务或等待重新授权",
            )
        if (
            requested_runtime is RuntimeVersion.PI
            and selected_runtime is not RuntimeVersion.PI
        ):
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="当前 Rollout 模式未向该用户开放 vNext",
            )
        payload.runtime_version = selected_runtime
        if (
            payload.table_output_contracts
            and selected_runtime is not RuntimeVersion.PI
        ):
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail="表格输出契约只能由 Pi Runtime 执行",
            )
    elif requested_runtime is RuntimeVersion.PI:
        # 显式请求 Pi 时不能绕过 Rollout/P0 门；路由状态不可读必须失败关闭。
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="Runtime 路由状态暂不可用，无法安全创建 Pi 任务",
        )
    else:
        payload.runtime_version = RuntimeVersion.LEGACY

    if (
        payload.table_output_contracts
        and payload.runtime_version is not RuntimeVersion.PI
    ):
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail="表格输出契约只能由 Pi Runtime 执行",
        )
    if (
        requested_runtime is None
        and payload.runtime_version is RuntimeVersion.LEGACY
    ):
        # 平台默认可能因 P0/路由状态落到 Legacy；不得遗留 Pi 连接、能力或外发确认。
        payload.model_connection_id = None
        payload.model_connection_model = None
        payload.capability_pack_refs = ()
        payload.validation_target = None
        payload.external_api_confirmed = False
        payload.provider = "local"
        payload.model = None
    capability_catalog = None
    if payload.capability_pack_refs:
        _require_capability_gray(user)
        if payload.runtime_version is not RuntimeVersion.PI:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail="任务能力只能冻结到 vNext TaskRevision",
            )
        capability_catalog = _capability_catalog()
        actor = catalog_actor_from_user(user)
        if payload.validation_target is not None:
            # #15 D9 验证目标资格：本人所有的个人包（平台包/他人包不能
            # 作为验证目标；gate 内仍强制其余三轴）。「必须同时被选择」
            # 已由模型 validator 保证。
            target_pack = capability_catalog.resolve_pack(
                actor,
                payload.validation_target.pack_id,
                payload.validation_target.version,
            )
            from src.conversation_steering import ProcedureScope

            if (
                target_pack is None
                or target_pack.digest != payload.validation_target.digest
                or target_pack.scope is not ProcedureScope.PERSONAL
                or target_pack.owner_id != user_id
            ):
                raise HTTPException(
                    status_code=status.HTTP_403_FORBIDDEN,
                    detail="验证目标必须是自己的个人能力",
                )
        # 创建与冻结共用装载门（AC4 同一公开 Interface）；digest 精确
        # 匹配仍由目录 freeze_selection 复核。
        _check_freeze_gate(
            actor,
            payload.capability_pack_refs,
            capability_catalog,
            validation_target=payload.validation_target,
        )
    connection_binding = None
    if (
        payload.runtime_version is RuntimeVersion.PI
        and payload.model_connection_id is None
    ):
        preference = get_default_broker().get_usage_preference(user_id)
        if preference is not None:
            if not preference["available"]:
                raise HTTPException(
                    status_code=status.HTTP_409_CONFLICT,
                    detail="默认模型连接已失效，请重新选择",
                )
            payload.model_connection_id = str(preference["connection_id"])
            payload.model_connection_model = str(preference["model_id"])
    if (
        payload.model_connection_id is not None
        and payload.runtime_version is not RuntimeVersion.PI
    ):
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail="模型连接当前只能用于 vNext 任务",
        )
    if payload.runtime_version is RuntimeVersion.PI:
        if not settings.pi_runtime_enabled:
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail="Pi 灰度入口当前未启用",
            )
        if payload.permission_profile is not PermissionProfile.STANDARD:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail=(
                    "当前只开放任务级容器标准增强模式；"
                    "扩展目录和宿主机权限需要单任务授权范围"
                ),
            )
        if payload.model_connection_id is not None:
            broker = get_default_broker()
            selected_connection = broker.get_connection(
                user_id,
                payload.model_connection_id,
            )
            if selected_connection is None:
                raise HTTPException(
                    status_code=status.HTTP_404_NOT_FOUND,
                    detail="模型连接不存在或无权访问",
                )
            if not payload.external_api_confirmed:
                raise HTTPException(
                    status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                    detail="使用外部模型连接前必须确认本任务的数据外发",
                )
            if (
                payload.model_connection_model is not None
                and not any(
                    item["model_id"] == payload.model_connection_model
                    and item["status"] == "available"
                    and item["enabled"]
                    for item in selected_connection.get("models", [])
                )
            ):
                raise HTTPException(
                    status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail="任务模型必须与已验证连接冻结的模型一致",
                )
            try:
                connection_binding = broker.freeze_connection(
                    user_id,
                    payload.model_connection_id,
                )
            except GrantError as exc:
                raise HTTPException(
                    status_code=status.HTTP_404_NOT_FOUND,
                    detail="模型连接不存在或无权访问",
                ) from exc
        elif not is_admin_role(user.get("role")):
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="普通用户需要选择自己的连接或管理员发布的连接",
            )
        elif payload.provider.lower() != "local":
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail="Pi 首期生产灰度只允许使用本地模型",
            )
    upload_store = _uploads()
    source_refs = []
    input_suffixes = set()
    for upload_id in payload.upload_ids:
        try:
            upload = upload_store.resolve(user_id, upload_id)
        except PermissionError as exc:
            raise HTTPException(
                status_code=404,
                detail="上传文件不存在或无权访问",
            ) from exc
        source_refs.append({
            "upload_id": upload.upload_id,
            "sha256": upload.sha256,
        })
        input_suffixes.add(Path(upload.original_name).suffix.lower())
    source_snapshot = None
    if payload.source_snapshot_id is not None:
        source_snapshot = SourceAcquisitionRepository(
            settings.webui_db_path
        ).get_snapshot(user_id, payload.source_snapshot_id)
        if source_snapshot is None:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="网页来源快照不存在或无权访问",
            )
        if int(source_snapshot["valid_page_count"]) < 1:
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="网页来源快照没有可执行的有效页面",
            )
        if source_snapshot.get("coverage", {}).get("status") == "hard_insufficient":
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="网页来源未达到启动前冻结的硬性有效页数；可刷新来源或调整要求",
            )
        if (
            source_snapshot.get("coverage", {}).get("status")
            == "coverage_unknown"
            and _HARD_SOURCE_REQUIREMENT_PATTERN.search(
                "\n".join((
                    payload.objective_text,
                    *payload.must_include,
                    payload.quantity_requirement or "",
                    payload.completeness_requirement or "",
                ))
            )
        ):
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail=(
                    "当前来源仍有失败、截断或未覆盖页面，不能承诺‘全部/所有/"
                    "必须 N 项’；请改为探索性要求，或重新获取完整授权范围"
                ),
            )
        source_refs = [
            {
                "kind": "web_artifact",
                "snapshot_id": payload.source_snapshot_id,
                "artifact_id": artifact["artifact_id"],
                "sha256": artifact["content_sha256"],
            }
            for artifact in source_snapshot["artifacts"]
        ]
        constraint_lines = [
            *(f"必须包含：{item}" for item in payload.must_include),
            *(f"明确不要：{item}" for item in payload.explicit_exclusions),
            f"数量要求：{payload.quantity_requirement}",
            f"完整性边界：{payload.completeness_requirement}",
        ]
        # 约束必须进入同一冻结 Revision，Agent 与独立 Verifier 才会执行，
        # 不能只把字段存进旁路合同而让运行时看不到。
        payload.objective_text = "\n".join(
            (user_objective, "", "执行约束：", *constraint_lines)
        )
    context_preview = None
    if payload.context_selection is not None:
        try:
            context_preview = _task_context_service().preview(
                owner_id=user_id,
                purpose=payload.context_purpose,
                objective_text=user_objective,
                output_formats=payload.output_formats,
                selection=payload.context_selection,
            )
        except KeyError as exc:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=str(exc),
            ) from exc
        except ValueError as exc:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail=str(exc),
            ) from exc
        if context_preview.preview_sha256 != payload.context_preview_sha256:
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="模板或记忆已变化，请重新检查并确认上下文草案",
            )
    if (
        payload.runtime_version is RuntimeVersion.LEGACY
        and input_suffixes & _DOCUMENT_INPUTS
        and input_suffixes & _TABLE_INPUTS
    ):
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail=(
                "当前版本暂不支持在同一任务中混合处理文档和表格；"
                "文档和表格请分别创建任务。"
            ),
        )
    requested_in_text = {
        "markdown" if item.lower() == "md" else item.lower()
        for item in _OUTPUT_FORMAT_PATTERN.findall(payload.objective_text)
    }
    selected_formats = set(payload.output_formats)
    missing_formats = requested_in_text - selected_formats
    if missing_formats:
        requested_label = "、".join(
            item.upper() for item in sorted(requested_in_text)
        )
        selected_label = "、".join(
            item.upper() for item in sorted(selected_formats)
        )
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail=(
                f"文字要求输出 {requested_label}，"
                f"但界面选择了 {selected_label}；请统一后再执行。"
            ),
        )
    repository = _runtime_repository()
    claimed_key = None
    if idempotency_key is not None:
        claimed_key = idempotency_key.strip()
        if not claimed_key or len(claimed_key) > 200:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail="Idempotency-Key 长度必须为 1 至 200 个字符",
            )
        fingerprint = hashlib.sha256(
            json.dumps(
                {
                    "payload": idempotency_payload,
                    "sources": source_refs,
                },
                ensure_ascii=False,
                sort_keys=True,
                separators=(",", ":"),
            ).encode("utf-8")
        ).hexdigest()
        try:
            task_id, is_new_claim = repository.claim_idempotency(
                user_id,
                claimed_key,
                request_hash=fingerprint,
                proposed_task_id=task_id,
            )
        except ValueError as exc:
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail=str(exc),
            ) from exc
        if not is_new_claim:
            existing = get_store().get_semantic_workspace_task(
                user_id,
                task_id,
            )
            if existing is None:
                raise HTTPException(
                    status_code=status.HTTP_409_CONFLICT,
                    detail="同一幂等请求正在创建，请稍后重试",
                )
            existing_runtime = repository.get(
                user_id,
                task_id,
                int(existing["active_revision"]),
            )
            return {
                **existing,
                "runtime_version": (
                    existing_runtime["runtime_version"].value
                    if existing_runtime
                    else RuntimeVersion.LEGACY.value
                ),
                "permission_profile": (
                    existing_runtime["permission_profile"].value
                    if existing_runtime
                    else PermissionProfile.STANDARD.value
                ),
                "model_connection_id": (
                    existing_runtime["model_connection_id"]
                    if existing_runtime
                    else None
                ),
            }
    transaction_hook = None
    if routing_repository is not None:
        assert routing_ref is not None
        assert rollout_actor is not None
        assert routing_preview is not None
    runtime_config = RuntimeTaskConfig(
        user_id=user_id,
        task_id=task_id,
        revision=1,
        runtime_version=payload.runtime_version,
        permission_profile=payload.permission_profile,
        model_connection_id=payload.model_connection_id,
        model_connection_version=(
            connection_binding.connection_version
            if connection_binding
            else None
        ),
        model_connection_model=(
            payload.model_connection_model or connection_binding.model
            if connection_binding
            else None
        ),
        external_api_confirmed=bool(connection_binding),
    )
    prepared_binding = None
    prepared_manifest = None
    if payload.source_snapshot_id is not None:
        try:
            prepared_binding, prepared_manifest = (
                await get_semantic_workspace_manager().prepare_runtime_binding(
                    model_connection_id=runtime_config.model_connection_id,
                    model_connection_version=(
                        runtime_config.model_connection_version
                    ),
                    model=(
                        runtime_config.model_connection_model
                        or payload.model
                        or settings.llm_model_name
                    ),
                )
            )
        except Exception:
            if claimed_key:
                repository.release_idempotency(
                    user_id,
                    claimed_key,
                    task_id=task_id,
                )
            raise
        runtime_config = runtime_config.model_copy(
            update={"run_id": prepared_binding.external_run_id}
        )
    payload.runtime_version, transaction_hook = _prepare_runtime_binding(
        _RevisionRoutingPlan(
            selected_runtime=payload.runtime_version,
            repository=routing_repository,
            task_revision=routing_ref,
            actor=rollout_actor,
            rollout=routing_preview,
        ),
        runtime_config,
    )
    if payload.source_snapshot_id is not None:
        assert prepared_binding is not None
        assert prepared_manifest is not None
        base_transaction_hook = transaction_hook
        assert source_snapshot is not None
        goal_contract = _freeze_goal_contract(
            payload,
            objective=user_objective,
            source_snapshot=source_snapshot,
        )
        delivery_spec = {"formats": list(payload.output_formats)}
        runtime_binding = prepared_binding.model_dump(mode="json")

        def bind_web_contract(connection: sqlite3.Connection) -> None:
            if base_transaction_hook is not None:
                base_transaction_hook(connection)
            _runtime_repository().freeze_runtime_binding(
                user_id,
                task_id,
                1,
                run_id=prepared_binding.external_run_id,
                binding=runtime_binding,
                capability_manifest=prepared_manifest.model_dump(mode="json"),
                adopted_existing_run=False,
                preallocated_run=True,
                connection=connection,
            )
            connection.execute(
                "INSERT INTO web_task_contracts "
                "(owner_id, task_id, revision, source_snapshot_id, "
                "goal_contract_json, delivery_spec_json, runtime_binding_json, "
                "created_at) VALUES (?, ?, 1, ?, ?, ?, ?, ?)",
                (
                    user_id,
                    task_id,
                    payload.source_snapshot_id,
                    json.dumps(goal_contract, ensure_ascii=False),
                    json.dumps(delivery_spec, ensure_ascii=False),
                    json.dumps(runtime_binding, ensure_ascii=False),
                    datetime.now(timezone.utc).isoformat(),
                ),
            )
            if context_preview is not None:
                _task_context_service().freeze(
                    connection,
                    owner_id=user_id,
                    task_id=task_id,
                    revision=1,
                    preview=context_preview,
                    expected_preview_sha256=payload.context_preview_sha256 or "",
                )

        transaction_hook = bind_web_contract
    first_line = payload.objective_text.splitlines()[0].strip()
    title = first_line[:40] + ("…" if len(first_line) > 40 else "")
    store = get_store()
    try:
        task = store.create_semantic_workspace_task(
            user_id,
            task_id=task_id,
            title=title or "未命名任务",
            objective_text=payload.objective_text,
            upload_ids=list(payload.upload_ids),
            output_formats=list(payload.output_formats),
            provider=payload.provider.lower(),
            model=payload.model,
            external_api_confirmed=payload.external_api_confirmed,
            source_refs=source_refs,
            table_output_contracts=[
                item.model_dump(mode="json")
                for item in payload.table_output_contracts
            ],
            transaction_hook=transaction_hook,
        )
        if capability_catalog is not None:
            capability_catalog.freeze_selection(
                catalog_actor_from_user(user),
                task_id=task_id,
                revision=1,
                pack_refs=payload.capability_pack_refs,
                validation_target=payload.validation_target,
            )
    except RuntimeError as exc:
        if claimed_key:
            repository.release_idempotency(
                user_id,
                claimed_key,
                task_id=task_id,
            )
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail=str(exc),
        ) from exc
    except Exception:
        if claimed_key:
            repository.release_idempotency(
                user_id,
                claimed_key,
                task_id=task_id,
            )
        raise
    store.append_semantic_workspace_event(
        user_id,
        task_id,
        stage="queued",
        event_type="task_created",
        summary=(
            f"任务已创建，共 {len(payload.upload_ids)} 个文件"
            if payload.upload_ids
            else f"任务已创建，共 {len(source_refs)} 个网页来源"
        ),
    )
    get_semantic_workspace_manager().enqueue(user_id, task_id)
    return {
        **task,
        "runtime_version": payload.runtime_version.value,
        "permission_profile": payload.permission_profile.value,
        "model_connection_id": payload.model_connection_id,
        "model_connection_model": (
            payload.model_connection_model or connection_binding.model
            if connection_binding
            else None
        ),
    }


@router.get("/tasks")
def list_tasks(
    task_status: str | None = Query(default=None, alias="status"),
    deleted: bool = False,
    limit: int = Query(default=100, ge=1, le=500),
    user=Depends(get_current_user),
):
    tasks = get_store().list_semantic_workspace_tasks(
        user["user_id"],
        status=task_status,
        deleted=deleted,
        limit=limit,
    )
    for task in tasks:
        runtime = _runtime_repository().get(
            user["user_id"],
            task["task_id"],
            task["active_revision"],
        )
        task["runtime_version"] = (
            runtime["runtime_version"].value
            if runtime
            else RuntimeVersion.LEGACY.value
        )
        task["permission_profile"] = (
            runtime["permission_profile"].value
            if runtime
            else PermissionProfile.STANDARD.value
        )
        task["model_connection_id"] = (
            runtime["model_connection_id"] if runtime else None
        )
        task["agentic_runtime_status"] = (
            runtime["status"].value if runtime else None
        )
    return tasks


@router.get("/tasks/{task_id}")
def get_task(
    task_id: str,
    revision: int | None = Query(default=None, ge=1),
    user=Depends(get_current_user),
):
    return _task_detail(
        user["user_id"],
        task_id,
        revision=revision,
        audience=(
            ProgressAudience.ADMIN
            if is_admin_role(user.get("role"))
            else ProgressAudience.USER
        ),
    )


@router.get("/tasks/{task_id}/events")
def get_task_events(
    task_id: str,
    after: int = Query(default=0, ge=0),
    user=Depends(get_current_user),
):
    _task_or_404(user["user_id"], task_id)
    return get_store().list_semantic_workspace_events(
        user["user_id"], task_id, after=after
    )


@router.get("/tasks/{task_id}/stream")
def stream_task(
    task_id: str,
    user=Depends(get_current_user),
):
    user_id = user["user_id"]
    _task_or_404(user_id, task_id)
    audience = (
        ProgressAudience.ADMIN
        if is_admin_role(user.get("role"))
        else ProgressAudience.USER
    )

    def safe_progress_event(
        task: dict[str, Any],
        event: dict[str, Any],
        *,
        harness: bool = False,
    ) -> dict[str, Any] | None:
        projection_task = {
            **task,
            "viewing_revision": int(task["active_revision"]),
            "events": [] if harness else [event],
            "harness_events": [event] if harness else [],
        }
        projected = ProgressProjection().project(
            _structured_progress_events(projection_task),
            audience=audience,
            task_status=task["status"],
        )
        if not projected.events:
            return None
        return projected.events[0].model_dump(mode="json")

    async def event_gen():
        store = get_store()
        existing_events = store.list_semantic_workspace_events(
            user_id, task_id
        )
        revision_boundaries = [
            event["sequence"]
            for event in existing_events
            if event["event_type"] == "revision_created"
        ]
        workspace_after = (
            revision_boundaries[-1] - 1 if revision_boundaries else 0
        )
        harness_after = 0
        last_status = ""
        while True:
            task = store.get_semantic_workspace_task(user_id, task_id)
            if task is None:
                yield {
                    "event": "error",
                    "data": json.dumps(
                        {"message": "任务不存在"}, ensure_ascii=False
                    ),
                }
                break
            for event in store.list_semantic_workspace_events(
                user_id, task_id, after=workspace_after
            ):
                workspace_after = max(workspace_after, event["sequence"])
                public_event = safe_progress_event(task, event)
                if public_event is None:
                    continue
                yield {
                    "id": event["event_id"],
                    "event": "progress",
                    "data": json.dumps(
                        public_event,
                        ensure_ascii=False,
                    ),
                }
            if task["run_id"]:
                for event in store.list_semantic_harness_events(
                    user_id, task["run_id"]
                ):
                    if event["sequence"] <= harness_after:
                        continue
                    harness_after = event["sequence"]
                    public_event = safe_progress_event(
                        task,
                        event,
                        harness=True,
                    )
                    if public_event is None:
                        continue
                    yield {
                        "id": event["event_id"],
                        "event": "progress",
                        "data": json.dumps(
                            public_event,
                            ensure_ascii=False,
                        ),
                    }
            if task["status"] != last_status:
                last_status = task["status"]
                yield {
                    "event": "status",
                    "data": json.dumps(
                        {
                            "task_id": task_id,
                            "status": task["status"],
                            "question": task["question"],
                            "error": task["error"],
                            "failure": task["failure"],
                        },
                        ensure_ascii=False,
                    ),
                }
            if task["status"] in _TERMINAL:
                yield {
                    "event": "done",
                    "data": json.dumps(
                        {"status": task["status"]},
                        ensure_ascii=False,
                    ),
                }
                break
            await asyncio.sleep(0.5)

    return EventSourceResponse(
        event_gen(),
        ping=15,
        headers={"Cache-Control": "no-cache"},
    )


@router.post("/tasks/{task_id}/turns")
async def steer_task(
    task_id: str,
    payload: WorkspaceTurnIn,
    idempotency_key: str | None = Header(
        default=None,
        alias="Idempotency-Key",
    ),
    user=Depends(get_current_user),
):
    """理解运行中追问；在用户确认前绝不修改活动 revision。"""

    user_id = user["user_id"]
    task = _task_or_404(user_id, task_id)
    events = get_store().list_semantic_workspace_events(user_id, task_id)
    request = SteeringRequest(
        owner_id=user_id,
        task_id=task_id,
        revision=int(task["active_revision"]),
        run_id=task.get("run_id"),
        text=payload.text,
        idempotency_key=idempotency_key,
        current_status=task["status"],
        status_summary=task.get("summary") or "",
        current_goal=task.get("objective_text") or "",
        event_summaries=tuple(
            str(event.get("summary") or "") for event in events[-8:]
        ),
        provider=task.get("provider") or "local",
        model=task.get("model"),
        external_api_confirmed=bool(
            task.get("external_api_confirmed")
            or payload.external_api_confirmed
        ),
    )
    try:
        service = ConversationSteering(
            _steering_repository(),
            build_context_rewriter(request),
        )
        result = await service.handle_turn(request)
    except ValueError as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail=str(exc),
        ) from exc

    # HTTP 幂等重放不能制造重复进度事件；结果 ID 是同一语义回合的稳定边界。
    if not any(
        event.get("details", {}).get("steering_result_id")
        == result.result_id
        for event in events
    ):
        event_type = {
            SteeringAction.ANSWER_ONLY: "followup.answered_without_change",
            SteeringAction.NORMALIZED_NO_MATERIAL_CHANGE: "context.rewrite_completed",
            SteeringAction.REVISION_PROPOSAL: "context.revision_proposed",
            SteeringAction.NEW_TASK_PROPOSAL: "context.new_task_proposed",
            SteeringAction.PERMISSION_REQUEST: "context.permission_required",
        }[result.action]
        current_stage = events[-1]["stage"] if events else "understand"
        get_store().append_semantic_workspace_event(
            user_id,
            task_id,
            stage=current_stage,
            event_type=event_type,
            summary=result.answer or result.acknowledgement,
            details={
                "steering_result_id": result.result_id,
                "proposal_id": result.proposal_id,
                "revision": result.revision,
                "run_id": result.run_id,
            },
        )
    return result.model_dump(mode="json")


@router.get("/tasks/{task_id}/turns")
def list_steering_turns(
    task_id: str,
    user=Depends(get_current_user),
):
    """恢复原话、语义差异和待确认草案，不把转写冒充用户原话。"""

    user_id = user["user_id"]
    task = _task_or_404(user_id, task_id)
    repository = _steering_repository()
    proposals = []
    for proposal in repository.list_proposals(user_id, task_id):
        if (
            proposal.status is RevisionProposalStatus.PENDING
            and proposal.base_revision != int(task["active_revision"])
        ):
            proposal = repository.update_proposal(
                proposal.model_copy(update={"status": RevisionProposalStatus.EXPIRED})
            )
        proposals.append(proposal.model_dump(mode="json"))
    turns = repository.list_turns(user_id, task_id)
    return {
        "turns": [turn.model_dump(mode="json") for turn in turns],
        "deltas": [
            delta.model_dump(mode="json")
            for turn in turns
            if (delta := repository.get_delta_for_turn(user_id, turn.turn_id))
            is not None
        ],
        "results": [
            result.model_dump(mode="json")
            for turn in turns
            if (result := repository.get_result_for_turn(user_id, turn.turn_id))
            is not None
        ],
        "proposals": proposals,
    }


@router.post("/tasks/{task_id}/revision-proposals/{proposal_id}/reject")
def reject_steering_revision(
    task_id: str,
    proposal_id: str,
    user=Depends(get_current_user),
):
    """拒绝草案只改变草案状态，绝不取消或改写当前 Run。"""

    user_id = user["user_id"]
    _task_or_404(user_id, task_id)
    repository = _steering_repository()
    proposal = repository.get_proposal(user_id, proposal_id)
    if proposal is None or proposal.task_id != task_id:
        raise HTTPException(status_code=404, detail="Revision 草案不存在或无权访问")
    if proposal.status is not RevisionProposalStatus.PENDING:
        raise HTTPException(status_code=409, detail="Revision 草案已处理")
    rejected = repository.update_proposal(
        proposal.model_copy(update={"status": RevisionProposalStatus.REJECTED})
    )
    return rejected.model_dump(mode="json")


def _apply_confirmed_steering_revision(
    user: dict[str, Any],
    task_id: str,
    decision_id: str,
    *,
    external_api_confirmed: bool,
) -> dict[str, Any]:
    """把已确认的结构化差异应用为新版本；旧版本和旧 Run 保持可追溯。"""

    user_id = user["user_id"]
    repository = _steering_repository()
    decision = repository.get_decision(user_id, decision_id)
    if decision is None or decision.task_id != task_id:
        raise HTTPException(status_code=404, detail="Revision 决策不存在或无权访问")
    if decision.status is not RevisionDecisionStatus.READY_TO_APPLY:
        raise HTTPException(status_code=409, detail="Revision 决策尚未到可应用状态")
    proposal = repository.get_proposal(user_id, decision.proposal_id)
    delta = repository.get_delta(user_id, proposal.delta_id) if proposal else None
    if proposal is None or delta is None:
        raise HTTPException(status_code=409, detail="Revision 草案数据不完整")
    task = _task_or_404(user_id, task_id)
    if int(task["active_revision"]) != decision.base_revision:
        raise HTTPException(status_code=409, detail="活动版本已变化，请重新确认修改")

    previous_runtime = _runtime_repository().get(
        user_id,
        task_id,
        decision.base_revision,
    )
    routing_plan = _preview_revision_runtime(
        user,
        task_id=task_id,
        revision=decision.base_revision + 1,
        requested_runtime=(
            previous_runtime["runtime_version"]
            if previous_runtime
            else RuntimeVersion.LEGACY
        ),
    )
    selected_runtime = routing_plan.selected_runtime
    if (
        selected_runtime is not RuntimeVersion.PI
        and task.get("table_output_contracts")
    ):
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail="表格输出契约只能由 Pi Runtime 执行",
        )
    connection_binding = None
    if (
        selected_runtime is RuntimeVersion.PI
        and previous_runtime
        and previous_runtime["model_connection_id"]
    ):
        if not external_api_confirmed:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail="新版本的数据外发范围必须再次确认",
            )
        try:
            connection_binding = get_default_broker().freeze_connection(
                user_id,
                previous_runtime["model_connection_id"],
            )
        except GrantError as exc:
            raise HTTPException(
                status_code=404,
                detail="模型连接不存在或无权访问",
            ) from exc

    formats = list(task["output_formats"])
    if delta.output_delta:
        normalized = [item.strip().lower() for item in delta.output_delta]
        invalid = set(normalized) - _FORMATS
        if invalid:
            raise HTTPException(
                status_code=422,
                detail=f"语义草案包含不支持的输出格式：{sorted(invalid)}",
            )
        formats = normalized
    objective = (
        f"{task['objective_text']}\n\n"
        "已确认的上下文变更（仅应用下列差异，其余语义继承）：\n"
        f"{delta.normalized_text}"
    )
    # 路由分配是新 Revision 的线性化点；所有可预见校验必须先完成。
    runtime_config = RuntimeTaskConfig(
        user_id=user_id,
        task_id=task_id,
        revision=decision.base_revision + 1,
        runtime_version=selected_runtime,
        permission_profile=(
            previous_runtime["permission_profile"]
            if previous_runtime
            else PermissionProfile.STANDARD
        ),
        model_connection_id=(
            previous_runtime["model_connection_id"]
            if selected_runtime is RuntimeVersion.PI and previous_runtime
            else None
        ),
        model_connection_version=(
            connection_binding.connection_version if connection_binding else None
        ),
        model_connection_model=(
            previous_runtime["model_connection_model"]
            if selected_runtime is RuntimeVersion.PI and previous_runtime
            else None
        ),
        external_api_confirmed=bool(connection_binding),
    )
    selected_runtime, transaction_hook = _prepare_runtime_binding(
        routing_plan,
        runtime_config,
    )
    transaction_hook = _inherit_task_context_hook(
        transaction_hook,
        owner_id=user_id,
        source_task_id=task_id,
        source_revision=decision.base_revision,
        target_task_id=task_id,
        target_revision=decision.base_revision + 1,
        objective_text=objective,
        output_formats=tuple(formats),
    )
    try:
        revision = get_store().create_semantic_workspace_revision(
            user_id,
            task_id,
            objective_text=objective,
            output_formats=formats,
            change_summary=delta.normalized_text,
            expected_revision=decision.base_revision + 1,
            transaction_hook=transaction_hook,
        )
    except RuntimeError as exc:
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail=str(exc),
        ) from exc
    if selected_runtime is RuntimeVersion.PI:
        _inherit_capability_selection(
            user_id,
            source_task_id=task_id,
            source_revision=proposal.base_revision,
            target_task_id=task_id,
            target_revision=int(revision["revision"]),
        )
    applied = repository.update_decision(
        decision.model_copy(
            update={
                "status": RevisionDecisionStatus.APPLIED,
                "updated_at": datetime.now().astimezone(),
            }
        )
    )
    get_store().append_semantic_workspace_event(
        user_id,
        task_id,
        stage="queued",
        event_type="revision_created",
        summary=f"已确认修改并创建结果版本 V{revision['revision']}",
        details={
            "proposal_id": proposal.proposal_id,
            "decision_id": decision.decision_id,
            "base_revision": decision.base_revision,
            "material_changes": list(proposal.material_changes),
        },
    )
    return {
        "decision": applied.model_dump(mode="json"),
        "revision": revision,
    }


@router.post(
    "/tasks/{task_id}/revision-proposals/{proposal_id}/decision",
    status_code=status.HTTP_202_ACCEPTED,
)
async def decide_steering_revision(
    task_id: str,
    proposal_id: str,
    payload: WorkspaceRevisionDecisionIn,
    user=Depends(get_current_user),
):
    user_id = user["user_id"]
    store = get_store()
    task = _task_or_404(user_id, task_id)
    repository = _steering_repository()
    proposal = repository.get_proposal(user_id, proposal_id)
    if proposal is None or proposal.task_id != task_id:
        raise HTTPException(status_code=404, detail="Revision 草案不存在或无权访问")
    if int(task["active_revision"]) != proposal.base_revision:
        raise HTTPException(status_code=409, detail="活动版本已变化，请重新确认修改")
    previous_runtime = _runtime_repository().get(
        user_id,
        task_id,
        proposal.base_revision,
    )
    # 外发确认必须发生在保存决策和取消 Run 之前；否则 422 响应仍可能破坏旧任务。
    if (
        previous_runtime
        and previous_runtime["model_connection_id"]
        and not payload.external_api_confirmed
    ):
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail="新版本或独立任务的数据外发范围必须再次确认",
        )
    service = ConversationSteering(_steering_repository(), None)
    try:
        decision = service.decide_proposal(
            user_id,
            proposal_id,
            payload.mode,
            external_api_confirmed=payload.external_api_confirmed,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    if payload.mode is RevisionSwitchMode.NEW_TASK:
        delta = repository.get_delta(user_id, proposal.delta_id) if proposal else None
        if proposal is None or delta is None:
            raise HTTPException(status_code=409, detail="Revision 草案数据不完整")
        formats = list(task["output_formats"])
        if delta.output_delta:
            formats = [item.strip().lower() for item in delta.output_delta]
            invalid = set(formats) - _FORMATS
            if invalid:
                raise HTTPException(
                    status_code=422,
                    detail=f"语义草案包含不支持的输出格式：{sorted(invalid)}",
                )
        new_task_id = f"workspace_{uuid.uuid4().hex[:16]}"
        routing_plan = _preview_revision_runtime(
            user,
            task_id=new_task_id,
            revision=1,
            requested_runtime=(
                previous_runtime["runtime_version"]
                if previous_runtime
                else RuntimeVersion.LEGACY
            ),
        )
        selected_runtime = routing_plan.selected_runtime
        inherited_contracts = [
            item
            for item in task.get("table_output_contracts", [])
            if item.get("format") in formats
        ]
        if (
            selected_runtime is not RuntimeVersion.PI
            and inherited_contracts
        ):
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail="表格输出契约只能由 Pi Runtime 执行",
            )
        connection_binding = None
        if (
            selected_runtime is RuntimeVersion.PI
            and previous_runtime
            and previous_runtime["model_connection_id"]
        ):
            if not payload.external_api_confirmed:
                raise HTTPException(
                    status_code=422,
                    detail="独立任务的数据外发范围必须单独确认",
                )
            try:
                connection_binding = get_default_broker().freeze_connection(
                    user_id,
                    previous_runtime["model_connection_id"],
                )
            except GrantError as exc:
                raise HTTPException(
                    status_code=404,
                    detail="模型连接不存在或无权访问",
                ) from exc
        # 新任务的格式、权限与外发确认通过后才冻结 Runtime 分配。
        runtime_config = RuntimeTaskConfig(
            user_id=user_id,
            task_id=new_task_id,
            revision=1,
            runtime_version=selected_runtime,
            permission_profile=(
                previous_runtime["permission_profile"]
                if previous_runtime
                else PermissionProfile.STANDARD
            ),
            model_connection_id=(
                previous_runtime["model_connection_id"]
                if selected_runtime is RuntimeVersion.PI and previous_runtime
                else None
            ),
            model_connection_version=(
                connection_binding.connection_version
                if connection_binding
                else None
            ),
            model_connection_model=(
                previous_runtime["model_connection_model"]
                if selected_runtime is RuntimeVersion.PI and previous_runtime
                else None
            ),
            external_api_confirmed=bool(connection_binding),
        )
        selected_runtime, transaction_hook = _prepare_runtime_binding(
            routing_plan,
            runtime_config,
        )
        new_objective = (
            f"{task['objective_text']}\n\n"
            "已确认的独立任务差异：\n"
            f"{delta.normalized_text}"
        )
        transaction_hook = _inherit_task_context_hook(
            transaction_hook,
            owner_id=user_id,
            source_task_id=task_id,
            source_revision=proposal.base_revision,
            target_task_id=new_task_id,
            target_revision=1,
            objective_text=new_objective,
            output_formats=tuple(formats),
        )
        try:
            new_task = get_store().create_semantic_workspace_task(
                user_id,
                task_id=new_task_id,
                title=f"{task['title']}（独立任务）",
                objective_text=new_objective,
                upload_ids=task["upload_ids"],
                output_formats=formats,
                provider=task["provider"],
                model=task["model"],
                external_api_confirmed=bool(connection_binding),
                table_output_contracts=inherited_contracts,
                transaction_hook=transaction_hook,
            )
        except RuntimeError as exc:
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail=str(exc),
            ) from exc
        if selected_runtime is RuntimeVersion.PI:
            _inherit_capability_selection(
                user_id,
                source_task_id=task_id,
                source_revision=proposal.base_revision,
                target_task_id=new_task_id,
                target_revision=1,
            )
        applied = repository.update_decision(
            decision.model_copy(
                update={
                    "status": RevisionDecisionStatus.APPLIED,
                    "updated_at": datetime.now().astimezone(),
                }
            )
        )
        get_store().append_semantic_workspace_event(
            user_id,
            new_task_id,
            stage="queued",
            event_type="task.created_from_revision_proposal",
            summary="已从确认的修改草案创建独立任务",
            details={
                "source_task_id": task_id,
                "proposal_id": proposal_id,
            },
        )
        get_semantic_workspace_manager().enqueue(user_id, new_task_id)
        return {
            "decision": applied.model_dump(mode="json"),
            "revision": None,
            "new_task": new_task,
        }
    if payload.mode is RevisionSwitchMode.AFTER_SAFE_POINT:
        get_store().append_semantic_workspace_event(
            user_id,
            task_id,
            stage="execute",
            event_type="revision.waiting_safe_point",
            summary="已确认修改，将在当前原子步骤结束后切换版本",
            details={"decision_id": decision.decision_id},
        )
        return {"decision": decision.model_dump(mode="json"), "revision": None}

    if task["status"] not in _TERMINAL:
        await get_semantic_workspace_manager().cancel(user_id, task_id)
    response = _apply_confirmed_steering_revision(
        user,
        task_id,
        decision.decision_id,
        external_api_confirmed=payload.external_api_confirmed,
    )
    get_semantic_workspace_manager().enqueue(user_id, task_id)
    return response


@router.post("/tasks/{task_id}/answer")
async def answer_task(
    task_id: str,
    payload: WorkspaceAnswerIn,
    user=Depends(get_current_user),
):
    try:
        return await get_semantic_workspace_manager().answer(
            user["user_id"], task_id, payload.answer.strip()
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    except ValueError as exc:
        raise HTTPException(status_code=409, detail=str(exc)) from exc


@router.post("/tasks/{task_id}/cancel")
async def cancel_task(
    task_id: str,
    user=Depends(get_current_user),
):
    try:
        return await get_semantic_workspace_manager().cancel(
            user["user_id"], task_id
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@router.post("/tasks/{task_id}/candidate-verification/retry")
async def retry_candidate_verification(
    task_id: str,
    user=Depends(get_current_user),
):
    """旧同步入口已退役，避免绕过外发确认与独立发布门。"""

    del task_id, user
    raise HTTPException(
        status_code=status.HTTP_410_GONE,
        detail=(
            "旧候选重试入口已退役；请先读取重验 Offer，再使用 "
            "candidate-verifications 创建追加式验证 Attempt"
        ),
    )


@router.post(
    "/tasks/{task_id}/candidate-verifications",
    status_code=status.HTTP_202_ACCEPTED,
)
async def request_candidate_reverification(
    task_id: str,
    payload: CandidateReverificationIn,
    idempotency_key: str = Header(
        min_length=1,
        max_length=240,
        alias="Idempotency-Key",
    ),
    user=Depends(get_current_user),
):
    """创建完整候选重验 Attempt；执行与正式发布均不在 HTTP 连接内完成。"""

    if (
        is_admin_role(user.get("role"))
        and get_store().get_semantic_workspace_task(user["user_id"], task_id)
        is None
    ):
        # 管理角色可以看跨 Owner 管理元数据，但不能代替 TaskOwner 签发重验权威。
        raise HTTPException(
            status_code=403,
            detail="候选重验与历史权威恢复只能由 TaskOwner 发起",
        )
    try:
        attempt = await get_semantic_workspace_manager().request_candidate_reverification(
            owner_id=user["user_id"],
            task_id=task_id,
            expected_revision=payload.expected_revision,
            expected_previous_attempt_id=payload.expected_previous_attempt_id,
            expected_candidate_set_hash=payload.expected_candidate_set_hash,
            expected_target_ruleset_hash=payload.expected_target_ruleset_hash,
            legacy_ruleset_unknown_acknowledged=(
                payload.legacy_ruleset_unknown_acknowledged is True
            ),
            authorization_text_version=payload.authorization_text_version,
            external_api_confirmed=payload.external_api_confirmed,
            accept_duplicate_provider_cost=payload.accept_duplicate_provider_cost,
            historical_authority_recovery=(
                payload.historical_authority_recovery
            ),
            idempotency_key=idempotency_key,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    except PermissionError as exc:
        raise HTTPException(status_code=403, detail=str(exc)) from exc
    except ReverificationContractError as exc:
        raise HTTPException(status_code=422, detail=str(exc)) from exc
    except (ReverificationUnavailableError, sqlite3.OperationalError) as exc:
        raise HTTPException(status_code=503, detail="候选重验服务暂时不可用") from exc
    except (RuntimeError, ValueError) as exc:
        raise HTTPException(status_code=409, detail=str(exc)) from exc
    audience = (
        ProgressAudience.ADMIN
        if is_admin_role(user.get("role"))
        else ProgressAudience.USER
    )
    return {
        "attempt_id": attempt.attempt_id,
        "task_id": attempt.task_id,
        "revision": attempt.revision,
        "run_id": attempt.run_id,
        "previous_attempt_id": attempt.previous_attempt_id,
        "status": attempt.status.value,
        "task": _task_detail(user["user_id"], task_id, audience=audience),
    }


@router.post(
    "/tasks/{task_id}/candidate-verifications/{attempt_id}/publish",
)
async def publish_candidate_verification(
    task_id: str,
    attempt_id: str,
    payload: CandidateVerificationPublishIn,
    idempotency_key: str = Header(
        min_length=1,
        max_length=240,
        alias="Idempotency-Key",
    ),
    user=Depends(get_current_user),
):
    """显式发布精确 passed Attempt；不会隐式重新执行 Pi 或 Provider。"""

    try:
        delivery = await get_semantic_workspace_manager().publish_candidate_verification(
            owner_id=user["user_id"],
            task_id=task_id,
            expected_revision=payload.expected_revision,
            attempt_id=attempt_id,
            idempotency_key=idempotency_key,
        )
    except (KeyError, PermissionError) as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    except sqlite3.OperationalError as exc:
        raise HTTPException(
            status_code=503,
            detail="候选发布服务暂时不可用",
        ) from exc
    except (RuntimeError, ValueError) as exc:
        raise HTTPException(status_code=409, detail=str(exc)) from exc
    return delivery.model_dump(mode="json", exclude={"user_id"})


@router.post(
    "/tasks/{task_id}/revisions",
    status_code=status.HTTP_202_ACCEPTED,
)
async def create_revision(
    task_id: str,
    payload: WorkspaceRevisionIn,
    user=Depends(get_current_user),
):
    user_id = user["user_id"]
    store = get_store()
    task = _task_or_404(user_id, task_id)
    if int(task["active_revision"]) != payload.expected_active_revision:
        # 风险确认只对用户刚看到的失败版本有效，不能被旧页面重复使用。
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail="活动版本已变化，请查看最新结果后再决定是否重新执行",
        )
    current_web_contract = store.get_web_task_contract(
        user_id,
        task_id,
        int(task["active_revision"]),
    )
    if payload.source_snapshot_id is not None and current_web_contract is None:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail="只有网页来源任务可以刷新来源快照",
        )
    effective_snapshot_id = (
        payload.source_snapshot_id
        or (
            current_web_contract["source_snapshot_id"]
            if current_web_contract is not None
            else None
        )
    )
    effective_source_refs = list(task.get("source_refs", []))
    effective_snapshot = None
    if effective_snapshot_id is not None:
        effective_snapshot = SourceAcquisitionRepository(
            settings.webui_db_path
        ).get_snapshot(user_id, effective_snapshot_id)
        if effective_snapshot is None:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="网页来源快照不存在或无权访问",
            )
        if int(effective_snapshot["valid_page_count"]) < 1:
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="网页来源快照没有可执行的有效页面",
            )
        if effective_snapshot.get("coverage", {}).get("status") == "hard_insufficient":
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="网页来源未达到启动前冻结的硬性有效页数",
            )
        effective_source_refs = [
            {
                "kind": "web_artifact",
                "snapshot_id": effective_snapshot_id,
                "artifact_id": artifact["artifact_id"],
                "sha256": artifact["content_sha256"],
            }
            for artifact in effective_snapshot["artifacts"]
        ]
    previous_runtime = _runtime_repository().get(
        user_id,
        task_id,
        int(task["active_revision"]),
    )
    accepted_assessment = None
    if payload.accepted_candidate_set_hash is not None:
        if (
            previous_runtime is None
            or previous_runtime.get("verified_candidate_set_hash")
            != payload.accepted_candidate_set_hash
        ):
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="Candidate 已变化，请查看最新部分结果后再确认",
            )
        accepted_assessment = _runtime_repository().get_candidate_coverage(
            user_id=user_id,
            task_id=task_id,
            revision=int(task["active_revision"]),
            candidate_set_hash=payload.accepted_candidate_set_hash,
        )
        if (
            accepted_assessment is None
            or not accepted_assessment.is_partial
            or accepted_assessment.actual_result_count
            != payload.accepted_result_count
        ):
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="当前 Candidate 没有可接受的冻结缺口结论",
            )
    expected_revision = int(task["active_revision"]) + 1
    formats = (
        list(payload.output_formats)
        if payload.output_formats is not None
        else task["output_formats"]
    )
    if (
        payload.table_output_contracts is not None
        and not {
            item.format for item in payload.table_output_contracts
        }.issubset(formats)
    ):
        # 继承旧格式时也必须现场核对，不能把错配契约写入不可变 Revision。
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail="表格输出契约必须绑定正式输出格式",
        )
    routing_plan = _preview_revision_runtime(
        user,
        task_id=task_id,
        revision=expected_revision,
        requested_runtime=(
            previous_runtime["runtime_version"]
            if previous_runtime
            else RuntimeVersion.LEGACY
        ),
    )
    selected_runtime = routing_plan.selected_runtime
    effective_contracts = (
        payload.table_output_contracts
        if payload.table_output_contracts is not None
        else task.get("table_output_contracts", [])
    )
    if (
        effective_contracts
        and selected_runtime is not RuntimeVersion.PI
    ):
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
            detail="表格输出契约只能由 Pi Runtime 执行",
        )
    connection_binding = None
    if (
        selected_runtime is RuntimeVersion.PI
        and previous_runtime
        and previous_runtime["model_connection_id"]
    ):
        if not payload.external_api_confirmed:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_CONTENT,
                detail="外部连接任务创建新版本前必须再次确认本版本的数据外发",
            )
        try:
            connection_binding = get_default_broker().freeze_connection(
                user_id,
                previous_runtime["model_connection_id"],
            )
        except GrantError as exc:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="模型连接不存在或无权访问",
            ) from exc
    if task["status"] not in {
        "completed",
        "candidate_ready",
        "failed",
        "cancelled",
    }:
        await get_semantic_workspace_manager().cancel(user_id, task_id)
        task = _task_or_404(user_id, task_id)
    if int(task["active_revision"]) + 1 != expected_revision:
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail="活动版本已变化，请重新提交 Revision",
        )
    objective = (
        task["objective_text"]
        if payload.source_snapshot_id is not None
        else f"{task['objective_text']}\n用户修改要求：{payload.instruction}"
    )
    # 取消旧 Run 和全部输入校验完成后，再用 Rollout CAS 冻结本版本。
    runtime_config = RuntimeTaskConfig(
        user_id=user_id,
        task_id=task_id,
        revision=expected_revision,
        runtime_version=selected_runtime,
        permission_profile=(
            previous_runtime["permission_profile"]
            if previous_runtime
            else PermissionProfile.STANDARD
        ),
        model_connection_id=(
            previous_runtime["model_connection_id"]
            if selected_runtime is RuntimeVersion.PI and previous_runtime
            else None
        ),
        model_connection_version=(
            connection_binding.connection_version if connection_binding else None
        ),
        model_connection_model=(
            previous_runtime["model_connection_model"]
            if selected_runtime is RuntimeVersion.PI and previous_runtime
            else None
        ),
        external_api_confirmed=bool(connection_binding),
    )
    prepared_binding = None
    prepared_manifest = None
    if effective_snapshot_id is not None:
        prepared_binding, prepared_manifest = (
            await get_semantic_workspace_manager().prepare_runtime_binding(
                model_connection_id=runtime_config.model_connection_id,
                model_connection_version=runtime_config.model_connection_version,
                model=(
                    runtime_config.model_connection_model
                    or task.get("model")
                    or settings.llm_model_name
                ),
            )
        )
        runtime_config = runtime_config.model_copy(
            update={"run_id": prepared_binding.external_run_id}
        )
    selected_runtime, transaction_hook = _prepare_runtime_binding(
        routing_plan,
        runtime_config,
    )
    if effective_snapshot_id is not None:
        assert prepared_binding is not None
        assert prepared_manifest is not None
        base_transaction_hook = transaction_hook
        previous_goal = (
            current_web_contract["goal_contract"]
            if current_web_contract is not None
            else {"objective": task["objective_text"]}
        )
        if accepted_assessment is not None:
            previous_goal = json.loads(
                json.dumps(previous_goal, ensure_ascii=False)
            )
            coverage_goal = previous_goal.setdefault("coverage", {})
            coverage_goal["target_result_count"] = (
                accepted_assessment.actual_result_count
            )
            coverage_goal["accepted_gap_from"] = {
                "revision": int(task["active_revision"]),
                "candidate_set_hash": payload.accepted_candidate_set_hash,
                "previous_target_result_count": (
                    accepted_assessment.target_result_count
                ),
                "accepted_result_count": (
                    accepted_assessment.actual_result_count
                ),
            }
            previous_goal["quantity_requirement"] = (
                f"接受当前有证据的 {accepted_assessment.actual_result_count} 项"
            )
        delivery_spec = {"formats": formats}
        runtime_binding = prepared_binding.model_dump(mode="json")

        def bind_web_revision(connection: sqlite3.Connection) -> None:
            if base_transaction_hook is not None:
                base_transaction_hook(connection)
            _runtime_repository().freeze_runtime_binding(
                user_id,
                task_id,
                expected_revision,
                run_id=prepared_binding.external_run_id,
                binding=runtime_binding,
                capability_manifest=prepared_manifest.model_dump(mode="json"),
                adopted_existing_run=False,
                preallocated_run=True,
                connection=connection,
            )
            connection.execute(
                "INSERT INTO web_task_contracts "
                "(owner_id, task_id, revision, source_snapshot_id, "
                "goal_contract_json, delivery_spec_json, runtime_binding_json, "
                "created_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                (
                    user_id,
                    task_id,
                    expected_revision,
                    effective_snapshot_id,
                    json.dumps(previous_goal, ensure_ascii=False),
                    json.dumps(delivery_spec, ensure_ascii=False),
                    json.dumps(runtime_binding, ensure_ascii=False),
                    datetime.now(timezone.utc).isoformat(),
                ),
            )

        transaction_hook = bind_web_revision
    transaction_hook = _inherit_task_context_hook(
        transaction_hook,
        owner_id=user_id,
        source_task_id=task_id,
        source_revision=int(task["active_revision"]),
        target_task_id=task_id,
        target_revision=expected_revision,
        objective_text=objective,
        output_formats=tuple(formats),
    )
    try:
        revision = store.create_semantic_workspace_revision(
            user_id,
            task_id,
            objective_text=objective,
            output_formats=formats,
            change_summary=payload.instruction,
            source_refs=effective_source_refs,
            table_output_contracts=(
                [
                    item.model_dump(mode="json")
                    for item in payload.table_output_contracts
                ]
                if payload.table_output_contracts is not None
                else None
            ),
            expected_revision=expected_revision,
            transaction_hook=transaction_hook,
        )
    except RuntimeError as exc:
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail=str(exc),
        ) from exc
    if selected_runtime is RuntimeVersion.PI:
        _inherit_capability_selection(
            user_id,
            source_task_id=task_id,
            source_revision=int(task["active_revision"]),
            target_task_id=task_id,
            target_revision=int(revision["revision"]),
        )
    store.append_semantic_workspace_event(
        user_id,
        task_id,
        stage="queued",
        event_type="revision_created",
        summary=f"已创建结果版本 V{revision['revision']}",
    )
    get_semantic_workspace_manager().enqueue(user_id, task_id)
    return revision


@router.post(
    "/tasks/{task_id}/candidate-gap-actions",
    status_code=status.HTTP_202_ACCEPTED,
)
async def decide_candidate_gap(
    task_id: str,
    payload: CandidateGapActionIn,
    idempotency_key: str = Header(
        min_length=1,
        max_length=200,
        alias="Idempotency-Key",
    ),
    user=Depends(get_current_user),
):
    """记录 Owner 的单一缺口动作；只有接受缺口会创建新 Revision。"""

    user_id = user["user_id"]
    task = _task_or_404(user_id, task_id)
    repository = _runtime_repository()
    runtime = repository.get(user_id, task_id, payload.expected_revision)
    if (
        runtime is None
        or runtime.get("verified_candidate_set_hash")
        != payload.expected_candidate_set_hash
    ):
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail="Candidate 已变化，请查看最新部分结果后再决定",
        )
    assessment = repository.get_candidate_coverage(
        user_id=user_id,
        task_id=task_id,
        revision=payload.expected_revision,
        candidate_set_hash=payload.expected_candidate_set_hash,
    )
    if assessment is None or not assessment.is_partial:
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail="当前 Candidate 没有待处理的冻结缺口",
        )
    if assessment.same_run_repair_allowed and payload.action == "accept_gap":
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail="确认漏提只能在当前 Run 内修复，不能通过降低目标接受",
        )
    if payload.action == "accept_gap" and assessment.actual_result_count < 1:
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail="当前没有可接受的有证据结果，请补充或刷新来源",
        )
    request_hash = hashlib.sha256(
        json.dumps(
            payload.model_dump(mode="json"),
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    ).hexdigest()
    try:
        action, claimed = repository.claim_gap_action(
            user_id=user_id,
            task_id=task_id,
            idempotency_key=idempotency_key,
            request_hash=request_hash,
            source_revision=payload.expected_revision,
            candidate_set_hash=payload.expected_candidate_set_hash,
            action=payload.action,
        )
    except ValueError as exc:
        raise HTTPException(status_code=409, detail=str(exc)) from exc
    if not claimed and action["status"] == "completed":
        return {
            "action": payload.action,
            "status": "completed",
            "source_revision": payload.expected_revision,
            "target_revision": action["target_revision"],
        }
    if (
        not claimed
        and action["status"] == "pending"
        and payload.action == "accept_gap"
    ):
        # 新 Revision 已写入、动作终态尚未来得及回写时，从冻结合同恢复同一幂等结果。
        recovery_revision = payload.expected_revision + 1
        recovery_contract = get_store().get_web_task_contract(
            user_id,
            task_id,
            recovery_revision,
        )
        accepted_from = (
            (recovery_contract or {}).get("goal_contract", {})
            .get("coverage", {})
            .get("accepted_gap_from", {})
        )
        if (
            accepted_from.get("revision") == payload.expected_revision
            and accepted_from.get("candidate_set_hash")
            == payload.expected_candidate_set_hash
        ):
            completed = repository.complete_gap_action(
                user_id=user_id,
                task_id=task_id,
                idempotency_key=idempotency_key,
                target_revision=recovery_revision,
            )
            return {
                "action": payload.action,
                "status": completed["status"],
                "source_revision": payload.expected_revision,
                "target_revision": recovery_revision,
            }
    if int(task["active_revision"]) != payload.expected_revision:
        if claimed:
            repository.complete_gap_action(
                user_id=user_id,
                task_id=task_id,
                idempotency_key=idempotency_key,
                status="rejected",
            )
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail="活动版本已变化，请查看最新结果后再决定",
        )

    if payload.action == "accept_gap":
        revision = await create_revision(
            task_id,
            WorkspaceRevisionIn(
                instruction=(
                    f"接受当前 {assessment.actual_result_count} 项有证据结果，"
                    "并据此调整本版本目标"
                ),
                external_api_confirmed=payload.external_api_confirmed,
                expected_active_revision=payload.expected_revision,
                accepted_candidate_set_hash=payload.expected_candidate_set_hash,
                accepted_result_count=assessment.actual_result_count,
            ),
            user=user,
        )
        completed = repository.complete_gap_action(
            user_id=user_id,
            task_id=task_id,
            idempotency_key=idempotency_key,
            target_revision=int(revision["revision"]),
        )
        return {
            "action": payload.action,
            "status": completed["status"],
            "source_revision": payload.expected_revision,
            "target_revision": int(revision["revision"]),
        }

    repository.complete_gap_action(
        user_id=user_id,
        task_id=task_id,
        idempotency_key=idempotency_key,
    )
    next_actions = {
        "reject_gap": "保留当前版本，不创建正式交付",
        "supplement_source": "请添加或选择新的获准来源后创建新版本",
        "refresh_source": "请使用获取最新网页动作重新读取原获准范围",
    }
    get_store().append_semantic_workspace_event(
        user_id,
        task_id,
        stage="verify",
        event_type=f"candidate_gap.{payload.action}",
        summary=next_actions[payload.action],
        details={
            "source_revision": payload.expected_revision,
            "candidate_set_hash": payload.expected_candidate_set_hash,
            "formal_delivery": False,
        },
    )
    return {
        "action": payload.action,
        "status": "completed",
        "source_revision": payload.expected_revision,
        "target_revision": None,
        "next_action": next_actions[payload.action],
    }


@router.post(
    "/tasks/{task_id}/source-refresh",
    status_code=status.HTTP_202_ACCEPTED,
)
async def refresh_task_source(
    task_id: str,
    payload: SourceRefreshIn,
    idempotency_key: str = Header(
        min_length=1,
        max_length=200,
        alias="Idempotency-Key",
    ),
    user=Depends(get_current_user),
):
    """按旧版本冻结范围获取新快照，成功后才切换到新 Revision。"""

    user_id = user["user_id"]
    store = get_store()
    task = _task_or_404(user_id, task_id)
    contract = store.get_web_task_contract(
        user_id,
        task_id,
        payload.expected_active_revision,
    )
    if contract is None:
        raise HTTPException(status_code=404, detail="当前版本没有可刷新的网页来源")
    repository = SourceAcquisitionRepository(settings.webui_db_path)
    old_snapshot = repository.get_snapshot(
        user_id,
        contract["source_snapshot_id"],
    )
    if old_snapshot is None:
        raise HTTPException(status_code=409, detail="当前版本的冻结来源快照已缺失")
    old_attempt = repository.get_attempt(
        user_id,
        old_snapshot["attempt_id"],
        include_snapshot=False,
    )
    if old_attempt is None:
        raise HTTPException(status_code=409, detail="当前版本的来源获取事实已缺失")
    scope = old_snapshot["allowed_scope"]
    completeness = scope.get("completeness", {})
    refresh_key = (
        f"refresh-{task_id}-"
        f"{hashlib.sha256(idempotency_key.encode('utf-8')).hexdigest()}"
    )
    existing_attempt = repository.get_by_idempotency_key(user_id, refresh_key)
    if (
        existing_attempt is None
        and int(task["active_revision"]) != payload.expected_active_revision
    ):
        # 新请求不能在过期页面背后读取外部站点；已有请求仍可用同一键恢复结果。
        raise HTTPException(
            status_code=409,
            detail="活动版本已变化，请查看最新版本后再刷新来源",
        )
    try:
        attempt = await _source_acquisition_service().acquire(
            owner_id=user_id,
            idempotency_key=refresh_key,
            request=SourceAcquisitionRequest(
                url=scope.get("normalized_url", old_attempt["normalized_url"]),
                purpose=old_attempt["purpose"],
                scope_kind=scope.get("kind", "current_page"),
                page_limit=int(scope.get("page_limit", 1)),
                completeness_mode=completeness.get("mode", "exploratory"),
                required_valid_pages=completeness.get("required_valid_pages"),
                request_context=(
                    f"source-refresh:{task_id}:revision:"
                    f"{payload.expected_active_revision}"
                ),
            ),
            resume_unknown=payload.resume_unknown,
        )
    except AcquisitionConflictError as exc:
        raise HTTPException(status_code=409, detail=str(exc)) from exc
    if attempt["status"] == "acquiring":
        # 结果未知时不切换版本；调用方用相同幂等键恢复该 Attempt。
        return {"status": "acquiring", "attempt": attempt, "revision": None}
    if attempt["status"] != "succeeded" or attempt.get("snapshot") is None:
        raise HTTPException(
            status_code=409,
            detail={
                "message": "来源刷新未形成有效快照，任务仍使用旧版本",
                "attempt": attempt,
            },
        )
    snapshot = attempt["snapshot"]
    if snapshot.get("coverage", {}).get("status") == "hard_insufficient":
        raise HTTPException(
            status_code=409,
            detail={
                "message": "刷新结果未达到原版本冻结的硬性有效页数，任务仍使用旧版本",
                "attempt": attempt,
            },
        )
    existing_revision = store.find_web_task_revision_by_snapshot(
        user_id,
        task_id,
        snapshot["snapshot_id"],
    )
    if existing_revision is not None:
        return {
            "status": "revision_created",
            "attempt": attempt,
            "revision": store.get_semantic_workspace_revision(
                user_id,
                task_id,
                existing_revision,
            ),
        }
    refresh_request_hash = hashlib.sha256(
        json.dumps(
            {
                "task_id": task_id,
                "expected_revision": payload.expected_active_revision,
                "attempt_id": attempt["attempt_id"],
                "snapshot_id": snapshot["snapshot_id"],
            },
            sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    ).hexdigest()
    try:
        refresh_intent, claimed_binding = store.claim_source_refresh_intent(
            user_id,
            task_id,
            idempotency_key,
            request_hash=refresh_request_hash,
            expected_revision=payload.expected_active_revision,
            attempt_id=attempt["attempt_id"],
            snapshot_id=snapshot["snapshot_id"],
            resume_unknown=payload.resume_unknown,
        )
    except ValueError as exc:
        raise HTTPException(status_code=409, detail=str(exc)) from exc
    if not claimed_binding:
        bound_revision = (
            int(refresh_intent["created_revision"])
            if refresh_intent.get("created_revision") is not None
            else store.find_web_task_revision_by_snapshot(
                user_id,
                task_id,
                snapshot["snapshot_id"],
            )
        )
        if bound_revision is None:
            return {"status": "acquiring", "attempt": attempt, "revision": None}
        return {
            "status": "revision_created",
            "attempt": attempt,
            "revision": store.get_semantic_workspace_revision(
                user_id,
                task_id,
                bound_revision,
            ),
        }
    if int(task["active_revision"]) != payload.expected_active_revision:
        store.finish_source_refresh_intent(
            user_id,
            task_id,
            idempotency_key,
            revision=None,
        )
        raise HTTPException(
            status_code=409,
            detail="来源已刷新但活动版本发生变化；请保留幂等键并重新确认",
        )
    try:
        revision = await create_revision(
            task_id,
            WorkspaceRevisionIn(
                instruction="按原授权范围刷新网页来源",
                external_api_confirmed=payload.external_api_confirmed,
                expected_active_revision=payload.expected_active_revision,
                source_snapshot_id=snapshot["snapshot_id"],
            ),
            user,
        )
    except HTTPException as exc:
        existing_revision = store.find_web_task_revision_by_snapshot(
            user_id,
            task_id,
            snapshot["snapshot_id"],
        )
        if exc.status_code != 409 or existing_revision is None:
            store.finish_source_refresh_intent(
                user_id,
                task_id,
                idempotency_key,
                revision=None,
            )
            raise
        revision = store.get_semantic_workspace_revision(
            user_id,
            task_id,
            existing_revision,
        )
    except Exception:
        store.finish_source_refresh_intent(
            user_id,
            task_id,
            idempotency_key,
            revision=None,
        )
        raise
    assert revision is not None
    store.finish_source_refresh_intent(
        user_id,
        task_id,
        idempotency_key,
        revision=int(revision["revision"]),
    )
    return {
        "status": "revision_created",
        "attempt": attempt,
        "revision": revision,
    }


@router.get("/tasks/{task_id}/candidates/{artifact_id}")
def download_candidate(
    task_id: str,
    artifact_id: str,
    revision: int | None = Query(default=None, ge=1),
    user=Depends(get_current_user),
):
    """下载明确标记为“未验证候选”的 Pi 结果。"""

    user_id = user["user_id"]
    task = _task_or_404(user_id, task_id)
    selected_revision = revision or int(task["active_revision"])
    runtime = _runtime_repository().get(
        user_id, task_id, selected_revision
    )
    if runtime is None:
        raise HTTPException(status_code=404, detail="候选文件不存在")
    candidate = next(
        (
            item
            for item in runtime["candidates"]
            if item.artifact_id == artifact_id
        ),
        None,
    )
    if candidate is None:
        raise HTTPException(status_code=404, detail="候选文件不存在")
    if not _candidate_download_allowed(
        runtime["verification"],
        candidate_format=candidate.format,
    ):
        raise HTTPException(
            status_code=409,
            detail="候选缺少可信来源或身份校验，已禁止下载",
        )
    path = candidate.host_path.resolve()
    workspace_root = Path(runtime["workspace_root"] or "").resolve()
    # 下载时再次检查路径归属，避免数据库或容器产物被篡改后越过任务目录。
    if (
        not path.is_file()
        or path.is_symlink()
        or workspace_root not in path.parents
    ):
        raise HTTPException(
            status_code=409,
            detail="候选文件已失效或未通过路径校验",
        )
    return FileResponse(
        path,
        filename=candidate.filename,
        media_type="application/octet-stream",
        headers={
            "X-Mangrove-Artifact-Status": "unverified-candidate",
        },
    )


@router.delete("/tasks/{task_id}")
async def move_to_recycle_bin(
    task_id: str,
    user=Depends(get_current_user),
):
    user_id = user["user_id"]
    task = _task_or_404(user_id, task_id)
    if task["status"] not in _TERMINAL:
        await get_semantic_workspace_manager().cancel(user_id, task_id)
    return get_store().soft_delete_semantic_workspace_task(
        user_id, task_id
    )


@router.post("/tasks/{task_id}/restore")
def restore_task(
    task_id: str,
    user=Depends(get_current_user),
):
    task = _task_or_404(user["user_id"], task_id)
    if task["deleted_at"] is None:
        raise HTTPException(status_code=409, detail="任务不在回收站")
    return get_store().restore_semantic_workspace_task(
        user["user_id"], task_id
    )


@router.delete("/tasks/{task_id}/permanent")
def permanently_delete_task(
    task_id: str,
    user=Depends(get_current_user),
):
    task = _task_or_404(user["user_id"], task_id)
    if task["deleted_at"] is None:
        raise HTTPException(
            status_code=409,
            detail="必须先把任务移入回收站",
        )
    if not get_store().purge_semantic_workspace_task(
        user["user_id"], task_id
    ):
        raise HTTPException(status_code=404, detail="任务不存在")
    return {"ok": True}


def _quote_identifier(value: str) -> str:
    return '"' + value.replace('"', '""') + '"'


def _table_preview(
    path: Path,
    *,
    lineage_path: Path | None,
    offset: int,
    limit: int,
    search: str,
    sort_by: str | None,
    sort_direction: Literal["asc", "desc"],
) -> dict[str, Any]:
    with duckdb.connect(":memory:") as conn:
        suffix = path.suffix.lower()
        relation_sql = {
            ".parquet": "read_parquet(?)",
            ".csv": "read_csv_auto(?, header=true)",
            ".jsonl": "read_json_auto(?, format='newline_delimited')",
        }.get(suffix)
        if relation_sql is None:
            raise ValueError("当前表格类型不支持在线预览")
        all_columns = [
            row[0]
            for row in conn.execute(
                f"DESCRIBE SELECT * FROM {relation_sql}", [str(path)]
            ).fetchall()
        ]
        columns = [
            column
            for column in all_columns
            if column != "__mg_output_record_id"
        ]
        where = ""
        params: list[Any] = [str(path)]
        if search:
            clauses = [
                f"CAST({_quote_identifier(column)} AS VARCHAR) ILIKE ?"
                for column in columns
            ]
            where = " WHERE " + " OR ".join(clauses)
            params.extend([f"%{search}%"] * len(columns))
        total = conn.execute(
            f"SELECT COUNT(*) FROM {relation_sql}{where}", params
        ).fetchone()[0]
        order = ""
        if sort_by:
            if sort_by not in columns:
                raise ValueError("排序字段不存在")
            order = (
                f" ORDER BY {_quote_identifier(sort_by)} "
                f"{sort_direction.upper()} NULLS LAST"
            )
        rows = conn.execute(
            "WITH indexed AS ("
            f"SELECT *, row_number() OVER () - 1 AS __mg_result_index "
            f"FROM {relation_sql}) "
            f"SELECT * FROM indexed{where}{order} LIMIT ? OFFSET ?",
            [*params, limit, offset],
        ).fetchall()
        result_columns = [*all_columns, "__mg_result_index"]
        records = [
            {
                column: value
                for column, value in zip(result_columns, row)
            }
            for row in rows
        ]
        if lineage_path and lineage_path.is_file() and records:
            output_ids = [
                str(record["__mg_output_record_id"])
                for record in records
                if record.get("__mg_output_record_id")
            ]
            placeholders = ", ".join("?" for _ in output_ids)
            lineage_rows = conn.execute(
                "SELECT output_record_id, artifact_id, table_ref, "
                "row_number, evidence_json FROM read_parquet(?) "
                f"WHERE output_record_id IN ({placeholders}) "
                "ORDER BY output_record_id, artifact_id, row_number",
                [str(lineage_path), *output_ids],
            ).fetchall()
            grouped: dict[str, list[dict[str, Any]]] = {}
            for output_id, artifact_id, table_ref, row_number, evidence in (
                lineage_rows
            ):
                grouped.setdefault(str(output_id), []).append(
                    {
                        "artifact_id": artifact_id,
                        "table_ref": table_ref,
                        "row_number": row_number,
                        "values": json.loads(evidence),
                    }
                )
            for record in records:
                record.pop("__mg_result_index", None)
                output_id = str(
                    record.pop("__mg_output_record_id", "")
                )
                record["__lineage"] = grouped.get(output_id, [])
        else:
            for record in records:
                record.pop("__mg_result_index", None)
                record.pop("__mg_output_record_id", None)
        return {
            "kind": "table",
            "columns": columns,
            "rows": records,
            "total": total,
            "offset": offset,
            "limit": limit,
        }


def _python_table_preview(
    rows: list[dict[str, Any]],
    *,
    columns: list[str],
    offset: int,
    limit: int,
    search: str,
    sort_by: str | None,
    sort_direction: Literal["asc", "desc"],
) -> dict[str, Any]:
    if search:
        keyword = search.casefold()
        rows = [
            row
            for row in rows
            if any(keyword in str(row.get(column, "")).casefold() for column in columns)
        ]
    if sort_by:
        if sort_by not in columns:
            raise ValueError("排序字段不存在")

        def sort_key(row: dict[str, Any]) -> tuple[int, int, Any]:
            value = row.get(sort_by)
            if value is None:
                return (1, 0, "")
            if isinstance(value, (int, float)) and not isinstance(value, bool):
                return (0, 0, value)
            return (0, 1, str(value).casefold())

        rows.sort(key=sort_key, reverse=sort_direction == "desc")
    return {
        "kind": "table",
        "columns": columns,
        "rows": rows[offset : offset + limit],
        "total": len(rows),
        "offset": offset,
        "limit": limit,
    }


def _xlsx_preview(
    path: Path,
    *,
    offset: int,
    limit: int,
    search: str,
    sort_by: str | None,
    sort_direction: Literal["asc", "desc"],
) -> dict[str, Any]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        columns = ["工作表"]
        rows: list[dict[str, Any]] = []
        for sheet in workbook.worksheets:
            iterator = sheet.iter_rows(values_only=True)
            raw_headers = next(iterator, ())
            headers: list[str] = []
            seen: dict[str, int] = {}
            for index, value in enumerate(raw_headers):
                base = str(value).strip() if value is not None else f"列{index + 1}"
                seen[base] = seen.get(base, 0) + 1
                headers.append(base if seen[base] == 1 else f"{base}_{seen[base]}")
            for header in headers:
                if header not in columns:
                    columns.append(header)
            for values in iterator:
                if not any(value is not None for value in values):
                    continue
                row = {"工作表": sheet.title}
                row.update(
                    {
                        header: value
                        for header, value in zip(headers, values)
                    }
                )
                rows.append(row)
    finally:
        workbook.close()
    return _python_table_preview(
        rows,
        columns=columns,
        offset=offset,
        limit=limit,
        search=search,
        sort_by=sort_by,
        sort_direction=sort_direction,
    )


def _paginate_document_items(
    items: list[dict[str, Any]],
    *,
    action: str,
    offset: int,
    limit: int,
    search: str,
) -> dict[str, Any]:
    if search:
        keyword = search.casefold()
        items = [
            item
            for item in items
            if keyword
            in f"{item.get('label', '')} {item.get('content', '')}".casefold()
        ]
    return {
        "kind": "document",
        "action": action,
        "items": items[offset : offset + limit],
        "total": len(items),
        "offset": offset,
        "limit": limit,
        "warnings": [],
    }


def _file_document_preview(
    path: Path,
    *,
    offset: int,
    limit: int,
    search: str,
) -> dict[str, Any]:
    suffix = path.suffix.lower()
    sections: list[tuple[str, str]] = []
    if suffix in {".txt", ".md", ".markdown"}:
        text = path.read_text(encoding="utf-8-sig")
        sections = [
            (f"段落 {index}", value.strip())
            for index, value in enumerate(re.split(r"\n\s*\n", text), 1)
            if value.strip()
        ]
    elif suffix in {".html", ".htm"}:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(path.read_text(encoding="utf-8-sig"), "html.parser")
        blocks = [
            value.strip()
            for value in soup.get_text("\n").splitlines()
            if value.strip()
        ]
        sections = [(f"内容 {index}", value) for index, value in enumerate(blocks, 1)]
    elif suffix == ".docx":
        from docx import Document

        document = Document(path)
        paragraphs = [item.text.strip() for item in document.paragraphs if item.text.strip()]
        sections.extend(
            (f"段落 {index}", value) for index, value in enumerate(paragraphs, 1)
        )
        for table_index, table in enumerate(document.tables, 1):
            content = "\n".join(
                " | ".join(cell.text.strip() for cell in row.cells)
                for row in table.rows
            ).strip()
            if content:
                sections.append((f"表格 {table_index}", content))
    elif suffix == ".pdf":
        from pypdf import PdfReader

        sections = [
            (f"第 {index} 页", (page.extract_text() or "").strip())
            for index, page in enumerate(PdfReader(path).pages, 1)
        ]
    elif suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(path)
        for index, slide in enumerate(presentation.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            sections.append((f"第 {index} 页", "\n".join(texts)))
    else:
        raise ValueError("当前文档类型不支持在线预览")
    items = [
        {
            "type": "derived",
            "id": f"file-{index}",
            "label": label,
            "content": content or "（该页没有可提取文本）",
            "evidence_refs": [],
        }
        for index, (label, content) in enumerate(sections, 1)
    ]
    return _paginate_document_items(
        items,
        action=f"{suffix.lstrip('.')}_preview",
        offset=offset,
        limit=limit,
        search=search,
    )


def _document_preview(
    path: Path,
    *,
    offset: int,
    limit: int,
    search: str,
) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    items: list[dict[str, Any]] = []
    document_payload = payload if isinstance(payload, dict) else {}
    for item in document_payload.get("passages", []):
        items.append(
            {
                "type": "passage",
                "id": item["passage_id"],
                "label": item["label"],
                "content": item["text"],
                "evidence_refs": item["evidence_refs"],
            }
        )
    for item in document_payload.get("differences", []):
        content = "\n".join(
            value
            for value in (item.get("before"), item.get("after"))
            if value
        )
        items.append(
            {
                "type": "difference",
                "id": item["diff_id"],
                "label": item["label"],
                "content": content,
                "change_type": item["change_type"],
                "evidence_refs": [
                    *item.get("before_evidence", []),
                    *item.get("after_evidence", []),
                ],
            }
        )
    for item in document_payload.get("findings", []):
        items.append(
            {
                "type": "finding",
                "id": item["finding_id"],
                "label": item["label"],
                "content": item["message"],
                "status": item["status"],
                "evidence_refs": item["evidence_refs"],
            }
        )
    for item in document_payload.get("derived_content", []):
        items.append(
            {
                "type": "derived",
                "id": item["content_id"],
                "label": item["action"],
                "content": item["content"],
                "evidence_refs": item["evidence_refs"],
            }
        )
    if not items:
        # Pi 可以按用户要求直接生成通用 JSON；它不是 Legacy DocumentExecutionResult，
        # 但仍应按稳定根键形成可读项，不能让合法正式交付显示为空白。
        values = (
            list(payload.items())
            if isinstance(payload, dict)
            else list(enumerate(payload, start=1))
            if isinstance(payload, list)
            else [("结果", payload)]
        )
        for index, (label, value) in enumerate(values, start=1):
            content = (
                value
                if isinstance(value, str)
                else json.dumps(
                    value,
                    ensure_ascii=False,
                    indent=2,
                    default=str,
                )
            )
            items.append(
                {
                    "type": "derived",
                    "id": f"json-{index}",
                    "label": str(label),
                    "content": content,
                    "evidence_refs": [],
                }
            )
    if search:
        keyword = search.casefold()
        items = [
            item
            for item in items
            if keyword
            in f"{item.get('label', '')} {item.get('content', '')}".casefold()
        ]
    return {
        "kind": "document",
        "action": document_payload.get("action") or "json_preview",
        "items": items[offset : offset + limit],
        "total": len(items),
        "offset": offset,
        "limit": limit,
        "warnings": document_payload.get("warnings", []),
    }


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _preview_result_file(
    result_path: Path,
    *,
    lineage_path: Path | None,
    offset: int,
    limit: int,
    search: str,
    sort_by: str | None,
    sort_direction: Literal["asc", "desc"],
) -> dict[str, Any]:
    suffix = result_path.suffix.lower()
    if suffix in {".parquet", ".csv", ".jsonl"}:
        return _table_preview(
            result_path,
            lineage_path=lineage_path,
            offset=offset,
            limit=limit,
            search=search,
            sort_by=sort_by,
            sort_direction=sort_direction,
        )
    if suffix == ".xlsx":
        return _xlsx_preview(
            result_path,
            offset=offset,
            limit=limit,
            search=search,
            sort_by=sort_by,
            sort_direction=sort_direction,
        )
    if suffix == ".json":
        return _document_preview(
            result_path,
            offset=offset,
            limit=limit,
            search=search,
        )
    if suffix in {
        ".docx",
        ".pdf",
        ".html",
        ".htm",
        ".md",
        ".markdown",
        ".txt",
        ".pptx",
    }:
        return _file_document_preview(
            result_path,
            offset=offset,
            limit=limit,
            search=search,
        )
    raise ValueError("当前结果类型不支持在线预览")


@router.get("/tasks/{task_id}/preview")
def preview_task_result(
    task_id: str,
    revision: int | None = Query(default=None, ge=1),
    offset: int = Query(default=0, ge=0),
    limit: int = Query(default=100, ge=1, le=500),
    search: str = Query(default="", max_length=200),
    sort_by: str | None = Query(default=None),
    sort_direction: Literal["asc", "desc"] = "asc",
    user=Depends(get_current_user),
):
    user_id = user["user_id"]
    task = _task_or_404(user_id, task_id)
    selected_revision = (
        get_store().get_semantic_workspace_revision(
            user_id, task_id, revision
        )
        if revision is not None
        else get_store().get_semantic_workspace_revision(
            user_id, task_id, task["active_revision"]
        )
    )
    if selected_revision is None:
        raise HTTPException(status_code=404, detail="结果版本不存在")
    run_id = selected_revision["run_id"]
    if not run_id:
        raise HTTPException(status_code=409, detail="任务尚无可预览结果")
    store = get_store()
    delivery = store.latest_semantic_delivery(user_id, run_id)
    formal_output = next(
        (
            output
            for output in (delivery or {}).get("outputs", [])
            if output.get("format") in _FORMATS
        ),
        None,
    )
    root = Path(settings.semantic_execution_root).resolve()
    paths = store.latest_semantic_harness_artifact_paths(user_id, run_id)
    legacy_result = Path(paths.get("result", "")).resolve()
    legacy_available = (
        legacy_result.is_file()
        and (legacy_result == root or root in legacy_result.parents)
    )
    lineage_value = paths.get("lineage") if legacy_available else None
    if legacy_available:
        # Legacy 的 Parquet 携带逐行来源；正式 Excel 只负责下载，不能反过来削弱预览证据。
        result_path = legacy_result
    elif formal_output is not None:
        record = store.get_semantic_delivery_output(
            user_id,
            formal_output["output_id"],
        )
        if record is None:
            raise HTTPException(
                status_code=409,
                detail="正式交付预览文件登记缺失",
            )
        result_path = Path(record["file_path"]).resolve()
        if (
            (result_path != root and root not in result_path.parents)
            or not result_path.is_file()
            or result_path.stat().st_size != record["size_bytes"]
            or _sha256_file(result_path) != record["sha256"]
        ):
            raise HTTPException(
                status_code=409,
                detail="正式交付预览文件完整性校验失败",
            )
    else:
        result_path = legacy_result
    lineage_path = (
        Path(lineage_value).resolve() if lineage_value else None
    )
    if lineage_path is not None and (
        not lineage_path.is_file()
        or lineage_path != root
        and root not in lineage_path.parents
    ):
        lineage_path = None
    if (
        not result_path.is_file()
        or result_path != root
        and root not in result_path.parents
    ):
        raise HTTPException(status_code=404, detail="预览制品不存在")
    try:
        return _preview_result_file(
            result_path,
            lineage_path=lineage_path,
            offset=offset,
            limit=limit,
            search=search,
            sort_by=sort_by,
            sort_direction=sort_direction,
        )
    except (ValueError, OSError, duckdb.Error, json.JSONDecodeError) as exc:
        raise HTTPException(status_code=409, detail=str(exc)) from exc


def _safe_archive_name(name: str) -> str:
    value = Path(name).name.strip()
    value = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", value)
    return value or "file"


@router.get("/tasks/{task_id}/bundle")
def download_bundle(
    task_id: str,
    include_sources: bool = False,
    revision: int | None = Query(default=None, ge=1),
    user=Depends(get_current_user),
):
    user_id = user["user_id"]
    task = _task_or_404(user_id, task_id)
    selected_revision = (
        get_store().get_semantic_workspace_revision(
            user_id, task_id, revision
        )
        if revision is not None
        else get_store().get_semantic_workspace_revision(
            user_id, task_id, task["active_revision"]
        )
    )
    if selected_revision is None:
        raise HTTPException(status_code=404, detail="结果版本不存在")
    run_id = selected_revision["run_id"]
    if not run_id:
        raise HTTPException(status_code=409, detail="任务尚无正式交付")
    store = get_store()
    manifest = store.latest_semantic_delivery(user_id, run_id)
    if manifest is None:
        raise HTTPException(status_code=404, detail="正式交付不存在")
    bundle_root = (
        Path(settings.semantic_execution_root) / "_bundles"
    ).resolve()
    bundle_root.mkdir(parents=True, exist_ok=True)
    bundle_path = bundle_root / f"{task_id}-{uuid.uuid4().hex[:8]}.zip"
    trace = {
        "task": {
            key: value
            for key, value in task.items()
            if key not in {"question"}
        },
        "workspace_events": store.list_semantic_workspace_events(
            user_id, task_id
        ),
        "harness_events": store.list_semantic_harness_events(
            user_id, run_id
        ),
        "attempts": store.list_semantic_harness_attempts(
            user_id, run_id
        ),
    }
    qa = {
        item["filename"]: item["qa"] for item in manifest["outputs"]
    }
    with zipfile.ZipFile(
        bundle_path, "w", compression=zipfile.ZIP_DEFLATED
    ) as archive:
        for output in manifest["outputs"]:
            record = store.get_semantic_delivery_output(
                user_id, output["output_id"]
            )
            if record is None:
                raise HTTPException(
                    status_code=409,
                    detail=f"交付文件登记缺失：{output['filename']}",
                )
            path = Path(record["file_path"]).resolve()
            if not path.is_file():
                raise HTTPException(
                    status_code=409,
                    detail=f"交付文件不存在：{output['filename']}",
                )
            archive.write(
                path,
                arcname=f"outputs/{_safe_archive_name(output['filename'])}",
            )
        archive.writestr(
            "manifest.json",
            json.dumps(manifest, ensure_ascii=False, indent=2),
        )
        archive.writestr(
            "qa.json",
            json.dumps(qa, ensure_ascii=False, indent=2),
        )
        archive.writestr(
            "trace.json",
            json.dumps(trace, ensure_ascii=False, indent=2),
        )
        if include_sources:
            for upload_id in task["upload_ids"]:
                item = _uploads().resolve(user_id, upload_id)
                archive.write(
                    Path(item.storage_path),
                    arcname=(
                        "sources/"
                        + _safe_archive_name(item.original_name)
                    ),
                )
    return FileResponse(
        bundle_path,
        media_type="application/zip",
        filename=f"{_safe_archive_name(task['title'])}.zip",
        background=BackgroundTask(bundle_path.unlink, missing_ok=True),
    )


@router.get("/storage")
def storage_usage(user=Depends(get_current_user)):
    user_id = user["user_id"]
    store = get_store()
    tasks = [
        *store.list_semantic_workspace_tasks(
            user_id, deleted=False, limit=500
        ),
        *store.list_semantic_workspace_tasks(
            user_id, deleted=True, limit=500
        ),
    ]
    upload_ids = {
        upload_id
        for task in tasks
        for upload_id in task["upload_ids"]
    }
    upload_bytes = 0
    for upload_id in upload_ids:
        try:
            upload_bytes += _uploads().resolve(
                user_id, upload_id
            ).size_bytes
        except PermissionError:
            pass
    output_ids: set[str] = set()
    output_bytes = 0
    for task in tasks:
        for revision in store.list_semantic_workspace_revisions(
            user_id, task["task_id"]
        ):
            if not revision["run_id"]:
                continue
            delivery = store.latest_semantic_delivery(
                user_id, revision["run_id"]
            )
            if not delivery:
                continue
            for output in delivery["outputs"]:
                if output["output_id"] in output_ids:
                    continue
                output_ids.add(output["output_id"])
                output_bytes += output["size_bytes"]
    return {
        "task_count": len(tasks),
        "recycle_bin_count": sum(
            1 for task in tasks if task["deleted_at"]
        ),
        "upload_bytes": upload_bytes,
        "delivery_bytes": output_bytes,
        "total_bytes": upload_bytes + output_bytes,
        "retention": "用户删除前永久保留；回收站保留 30 天",
        "calculated_at": datetime.now().isoformat(timespec="seconds"),
    }
